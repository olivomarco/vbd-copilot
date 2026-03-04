#!/usr/bin/env python3
"""
Programmatic QA checks for generated .pptx files.

Runs a battery of layout, content, and consistency checks that catch
the most common generation issues WITHOUT requiring visual inspection.
Returns a structured JSON report with severity-tagged findings.

Usage:
    python scripts/pptx_qa_checks.py <path-to-pptx> [--expected-slides N]

Exit codes:
    0 = CLEAN (no CRITICAL or MAJOR issues)
    1 = ISSUES_FOUND (at least one CRITICAL or MAJOR)
    2 = ERROR (could not open file)
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches, Pt

# ── Layout constants (must match pptx_utils.py) ──────────────────────────────
SLIDE_WIDTH_EMU = Inches(13.333)
SLIDE_HEIGHT_EMU = Inches(7.5)
CONTENT_LEFT_EMU = Inches(0.8)
CONTENT_TOP_EMU = Inches(1.2)
CONTENT_WIDTH_EMU = Inches(11.0)
CONTENT_BOTTOM_EMU = Inches(6.8)

# Tolerances
OVERFLOW_TOLERANCE_EMU = Inches(0.15)  # shapes may exceed by up to 0.15" (decorative bleed)
MIN_MARGIN_EMU = Inches(0.3)          # minimum margin from any slide edge for content shapes
OVERLAP_THRESHOLD_EMU = Inches(0.5)   # shapes overlapping by more than this are flagged

# Font sanity range
MIN_FONT_PT = 8
MAX_FONT_PT = 48

# Placeholder patterns
PLACEHOLDER_RE = re.compile(
    r"xxxx|lorem\s+ipsum|placeholder|TODO|FIXME|insert\s+here|TBD|sample\s+text",
    re.IGNORECASE,
)


def emu_to_inches(emu: int) -> float:
    return round(emu / 914400, 2)


# ── Individual check functions ────────────────────────────────────────────────

def check_slide_count(prs: Presentation, expected: int | None) -> list[dict]:
    """Check actual vs expected slide count."""
    issues = []
    actual = len(prs.slides)
    if expected is not None and actual != expected:
        issues.append({
            "slide": 0,
            "severity": "CRITICAL" if abs(actual - expected) > 2 else "MAJOR",
            "check": "slide_count",
            "message": f"Expected {expected} slides, got {actual}",
        })
    if actual == 0:
        issues.append({
            "slide": 0,
            "severity": "CRITICAL",
            "check": "slide_count",
            "message": "Presentation has zero slides",
        })
    return issues


def check_shape_overflow(prs: Presentation) -> list[dict]:
    """Detect shapes that extend beyond slide boundaries.

    Properly handles grouped shapes by computing absolute positions
    through the group hierarchy.
    """
    issues = []
    sw = SLIDE_WIDTH_EMU
    sh = SLIDE_HEIGHT_EMU
    tol = OVERFLOW_TOLERANCE_EMU

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            # Skip grouped shapes - they're checked via their parent group
            if hasattr(shape, "_element"):
                parent = shape._element.getparent()
                if parent is not None and parent.tag.endswith("}grpSp"):
                    continue

            _check_shape_bounds(shape, slide_idx, sw, sh, tol, issues, 0, 0)

    return issues


def _check_shape_bounds(
    shape, slide_idx: int, sw: int, sh: int, tol: int,
    issues: list[dict], offset_x: int, offset_y: int
):
    """Recursively check shape bounds, accounting for group offsets."""
    if hasattr(shape, "shapes"):
        # This is a group shape - recurse into children with group offset
        grp_x = (shape.left or 0) + offset_x
        grp_y = (shape.top or 0) + offset_y
        for child in shape.shapes:
            _check_shape_bounds(child, slide_idx, sw, sh, tol, issues, grp_x, grp_y)
        return

    left = (shape.left or 0) + offset_x
    top = (shape.top or 0) + offset_y
    width = shape.width or 0
    height = shape.height or 0
    right = left + width
    bottom = top + height

    # Get shape name/text for reporting
    shape_name = getattr(shape, "name", "") or ""
    shape_text = ""
    if shape.has_text_frame:
        shape_text = shape.text_frame.text[:60]

    label = shape_name or shape_text or f"shape@({emu_to_inches(left)},{emu_to_inches(top)})"

    # Skip known decorative elements that intentionally bleed off-slide:
    # - Ovals on slide 1 (lead slide's concentric circle badge)
    # - Ovals on the closing slide (decorative background circle)
    name_lower = shape_name.lower()
    is_decorative_oval = "oval" in name_lower and not shape_text.strip()
    if is_decorative_oval:
        return  # decorative circles intentionally overflow

    # Check right overflow
    if right > sw + tol:
        overflow_in = emu_to_inches(right - sw)
        issues.append({
            "slide": slide_idx,
            "severity": "CRITICAL" if overflow_in > 0.5 else "MAJOR",
            "check": "shape_overflow_right",
            "message": (
                f"Shape '{label}' overflows right edge by {overflow_in}\" "
                f"(right={emu_to_inches(right)}\", slide_w={emu_to_inches(sw)}\")"
            ),
        })

    # Check bottom overflow
    if bottom > sh + tol:
        overflow_in = emu_to_inches(bottom - sh)
        issues.append({
            "slide": slide_idx,
            "severity": "CRITICAL" if overflow_in > 0.5 else "MAJOR",
            "check": "shape_overflow_bottom",
            "message": (
                f"Shape '{label}' overflows bottom edge by {overflow_in}\" "
                f"(bottom={emu_to_inches(bottom)}\", slide_h={emu_to_inches(sh)}\")"
            ),
        })

    # Check left/top negative overflow (off-screen left/top)
    if left < -tol and width > Inches(0.5):
        issues.append({
            "slide": slide_idx,
            "severity": "MAJOR",
            "check": "shape_overflow_left",
            "message": f"Shape '{label}' extends {emu_to_inches(-left)}\" off left edge",
        })


def check_empty_text_frames(prs: Presentation) -> list[dict]:
    """Find text frames that are likely intended for content but are empty.

    Excludes shapes that are clearly decorative/structural:
    - Background fills (Rounded Rectangle, Rectangle used as card backgrounds)
    - Thin accent bars/dividers (height < 0.1" or aspect ratio > 15:1)
    - Ovals/circles (decorative elements)
    - Shapes named with 'Rectangle' or 'Rounded Rectangle' without 'TextBox' in name
    """
    issues = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        text_shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                text_shapes.append((shape, text))

        has_content = any(len(t) > 0 for _, t in text_shapes)
        if has_content:
            for shape, text in text_shapes:
                if text:
                    continue
                w = shape.width or 0
                h = shape.height or 0
                name = getattr(shape, "name", "") or "unnamed"

                # Skip small shapes
                if w < Inches(1.5) or h < Inches(0.3):
                    continue
                # Skip thin bars/dividers
                if h < Inches(0.1):
                    continue
                if h > 0 and w / h > 15:
                    continue
                # Skip known decorative shape types (card backgrounds, ovals)
                name_lower = name.lower()
                if any(kw in name_lower for kw in ["rectangle", "oval", "freeform"]):
                    continue
                # Only flag TextBox-named shapes that are empty
                if "textbox" not in name_lower:
                    continue

                issues.append({
                    "slide": slide_idx,
                    "severity": "MINOR",
                    "check": "empty_text_frame",
                    "message": (
                        f"Empty text frame '{name}' "
                        f"({emu_to_inches(w)}\" x {emu_to_inches(h)}\")"
                    ),
                })
    return issues


def check_placeholder_text(prs: Presentation) -> list[dict]:
    """Scan all text for placeholder/template content."""
    issues = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                matches = PLACEHOLDER_RE.findall(text)
                if matches:
                    issues.append({
                        "slide": slide_idx,
                        "severity": "CRITICAL",
                        "check": "placeholder_text",
                        "message": (
                            f"Placeholder text found: {matches[:3]} "
                            f"in '{text[:80]}...'"
                        ),
                    })

        # Also check speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            matches = PLACEHOLDER_RE.findall(notes_text)
            if matches:
                issues.append({
                    "slide": slide_idx,
                    "severity": "MAJOR",
                    "check": "placeholder_in_notes",
                    "message": f"Placeholder text in speaker notes: {matches[:3]}",
                })
    return issues


def check_speaker_notes(prs: Presentation) -> list[dict]:
    """Verify speaker notes are present and substantial."""
    issues = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        has_notes = False
        notes_len = 0
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            has_notes = len(notes_text) > 0
            notes_len = len(notes_text)

        if not has_notes:
            issues.append({
                "slide": slide_idx,
                "severity": "MAJOR",
                "check": "missing_notes",
                "message": "No speaker notes on this slide",
            })
        elif notes_len < 50:
            issues.append({
                "slide": slide_idx,
                "severity": "MINOR",
                "check": "short_notes",
                "message": (
                    f"Speaker notes are very short ({notes_len} chars) - "
                    f"expected full presenter transcript"
                ),
            })
    return issues


def check_font_sizes(prs: Presentation) -> list[dict]:
    """Check for out-of-range font sizes."""
    issues = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None:
                        pt = run.font.size.pt
                        if pt < MIN_FONT_PT:
                            issues.append({
                                "slide": slide_idx,
                                "severity": "MAJOR",
                                "check": "font_too_small",
                                "message": (
                                    f"Font size {pt}pt is below minimum ({MIN_FONT_PT}pt): "
                                    f"'{run.text[:40]}'"
                                ),
                            })
                        elif pt > MAX_FONT_PT:
                            issues.append({
                                "slide": slide_idx,
                                "severity": "MINOR",
                                "check": "font_too_large",
                                "message": (
                                    f"Font size {pt}pt exceeds maximum ({MAX_FONT_PT}pt): "
                                    f"'{run.text[:40]}'"
                                ),
                            })
    return issues


def check_text_overflow_heuristic(prs: Presentation) -> list[dict]:
    """Heuristic check for text that likely overflows its bounding box.

    Estimates the number of lines based on character count and box width,
    then checks if that exceeds the available height. This is approximate
    because we don't have full font metrics, but catches obvious cases.
    """
    issues = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            tf = shape.text_frame
            box_w = shape.width or 0
            box_h = shape.height or 0
            if box_w < Inches(0.5) or box_h < Inches(0.3):
                continue  # skip tiny shapes

            # Estimate characters per line based on average char width
            # Segoe UI 14pt ~ 7pt avg char width ~ 0.097"
            # For safety, use a conservative estimate
            total_text = tf.text
            if not total_text.strip():
                continue

            # Get the dominant font size in this text frame
            font_sizes = []
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        font_sizes.append(run.font.size.pt)
            avg_font = sum(font_sizes) / len(font_sizes) if font_sizes else 14

            # Approximate chars per line and line height
            avg_char_width_in = avg_font * 0.007  # rough: 1pt font ~ 0.007" avg char width
            chars_per_line = max(1, int(emu_to_inches(box_w) / avg_char_width_in))
            line_height_in = avg_font * 1.5 / 72  # 1.5x line spacing

            # Count effective lines across all paragraphs
            total_lines = 0
            for para in tf.paragraphs:
                para_text = para.text
                if not para_text:
                    total_lines += 1  # empty para = one line
                else:
                    total_lines += max(1, len(para_text) / chars_per_line)

            estimated_height = total_lines * line_height_in
            available_height = emu_to_inches(box_h)

            if estimated_height > available_height * 1.3:  # 30% overflow threshold
                overflow_pct = int((estimated_height / available_height - 1) * 100)
                shape_name = getattr(shape, "name", "") or "unnamed"
                issues.append({
                    "slide": slide_idx,
                    "severity": "MAJOR" if overflow_pct > 50 else "MINOR",
                    "check": "text_overflow_heuristic",
                    "message": (
                        f"Text likely overflows '{shape_name}' by ~{overflow_pct}% "
                        f"(~{int(total_lines)} lines in {available_height}\" height, "
                        f"font ~{avg_font:.0f}pt)"
                    ),
                })
    return issues


def check_shape_overlap(prs: Presentation) -> list[dict]:
    """Detect content shapes that unintentionally overlap each other.

    In pptx_utils, cards are built with a background shape (Rounded Rectangle)
    plus TextBox children placed on top. This is INTENTIONAL layering.
    We only flag overlap when:
      - Both shapes have substantial text content (two text-bearing shapes colliding)
      - OR shapes of the same type overlap (e.g., two cards stacked wrongly)
    """
    issues = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        content_shapes = []
        for shape in slide.shapes:
            left = shape.left or 0
            top = shape.top or 0
            w = shape.width or 0
            h = shape.height or 0

            # Skip full-slide backgrounds and tiny shapes
            if w >= SLIDE_WIDTH_EMU * 0.9 and h >= SLIDE_HEIGHT_EMU * 0.9:
                continue
            if w < Inches(0.5) or h < Inches(0.3):
                continue
            # Skip thin decorative bars/dividers (aspect ratio > 20:1)
            if h > 0 and w / h > 20:
                continue
            if w > 0 and h / w > 20:
                continue

            has_text = shape.has_text_frame and shape.text_frame.text.strip()
            shape_type_name = getattr(shape, "name", "") or ""
            is_background = not has_text and (
                "Rectangle" in shape_type_name or "Oval" in shape_type_name
            )

            content_shapes.append({
                "left": left, "top": top, "right": left + w,
                "bottom": top + h, "name": shape_type_name,
                "text": shape.text_frame.text[:30].strip() if shape.has_text_frame else "",
                "has_text": bool(has_text),
                "is_background": is_background,
            })

        # O(n^2) overlap check - fine for typical slide shape counts (< 30)
        for i, a in enumerate(content_shapes):
            for b in content_shapes[i + 1:]:
                # Skip intentional layering: background card + text child on top
                if a["is_background"] and b["has_text"] and not a["has_text"]:
                    continue
                if b["is_background"] and a["has_text"] and not b["has_text"]:
                    continue
                # Skip when neither has text (two background/decorative shapes)
                if not a["has_text"] and not b["has_text"]:
                    continue

                # Calculate overlap rectangle
                ox = max(0, min(a["right"], b["right"]) - max(a["left"], b["left"]))
                oy = max(0, min(a["bottom"], b["bottom"]) - max(a["top"], b["top"]))

                if ox > OVERLAP_THRESHOLD_EMU and oy > OVERLAP_THRESHOLD_EMU:
                    overlap_area = emu_to_inches(ox) * emu_to_inches(oy)
                    a_area = emu_to_inches(a["right"] - a["left"]) * emu_to_inches(a["bottom"] - a["top"])
                    b_area = emu_to_inches(b["right"] - b["left"]) * emu_to_inches(b["bottom"] - b["top"])
                    min_area = min(a_area, b_area)

                    if min_area > 0 and overlap_area / min_area > 0.3:
                        a_label = a["text"] or a["name"] or "shape"
                        b_label = b["text"] or b["name"] or "shape"
                        issues.append({
                            "slide": slide_idx,
                            "severity": "MAJOR",
                            "check": "shape_overlap",
                            "message": (
                                f"Text-bearing shapes overlap: '{a_label}' and '{b_label}' "
                                f"by {overlap_area:.1f} sq in "
                                f"({int(overlap_area / min_area * 100)}% of smaller shape)"
                            ),
                        })
    return issues


def check_closing_slide_cta_count(prs: Presentation) -> list[dict]:
    """Validate that the closing slide doesn't have too many CTA cards.

    The create_closing_slide template supports max 3 CTA items before
    they overflow the right edge.
    """
    issues = []
    if len(prs.slides) < 1:
        return issues

    last_slide = prs.slides[-1]
    # Heuristic: closing slide has shapes with "cta" or link-like text
    cta_shapes = []
    for shape in last_slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.lower()
            if any(kw in text for kw in ["http", "learn more", "get started",
                                          "aka.ms", "github.com", "learn.microsoft"]):
                left = shape.left or 0
                if left > Inches(5.0):  # CTA area is in the right half
                    cta_shapes.append(shape)

    # Count distinct CTA groups (shapes at similar x positions)
    if len(cta_shapes) > 6:  # 2 shapes per CTA card (icon + text) x 3 = 6
        issues.append({
            "slide": len(prs.slides),
            "severity": "MAJOR",
            "check": "closing_cta_overflow",
            "message": (
                f"Closing slide appears to have too many CTA items "
                f"({len(cta_shapes)} link shapes detected, max recommended is 3 CTA cards)"
            ),
        })
    return issues


def check_content_margins(prs: Presentation) -> list[dict]:
    """Check that content shapes respect minimum margins from slide edges."""
    issues = []
    min_margin = MIN_MARGIN_EMU

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            left = shape.left or 0
            top = shape.top or 0
            w = shape.width or 0
            h = shape.height or 0

            # Skip background/decorative shapes
            if w >= SLIDE_WIDTH_EMU * 0.9:
                continue
            if w < Inches(0.5):
                continue
            if not shape.has_text_frame:
                continue
            if not shape.text_frame.text.strip():
                continue

            # Check margins - only flag if shape has substantial content
            text_len = len(shape.text_frame.text.strip())
            if text_len < 5:
                continue

            if left < min_margin and left > 0:
                issues.append({
                    "slide": slide_idx,
                    "severity": "MINOR",
                    "check": "margin_too_small_left",
                    "message": (
                        f"Shape '{shape.text_frame.text[:30]}' has only "
                        f"{emu_to_inches(left)}\" left margin (min: {emu_to_inches(min_margin)}\")"
                    ),
                })

            right = left + w
            right_margin = SLIDE_WIDTH_EMU - right
            if 0 < right_margin < min_margin:
                issues.append({
                    "slide": slide_idx,
                    "severity": "MINOR",
                    "check": "margin_too_small_right",
                    "message": (
                        f"Shape '{shape.text_frame.text[:30]}' has only "
                        f"{emu_to_inches(right_margin)}\" right margin"
                    ),
                })

    return issues


def check_slide_text_density(prs: Presentation) -> list[dict]:
    """Flag slides that are too text-heavy (wall of text) or completely empty."""
    issues = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        total_chars = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                total_chars += len(shape.text_frame.text)

        if total_chars == 0 and slide_idx > 1:
            # Allow slide 1 (title) to have minimal text
            issues.append({
                "slide": slide_idx,
                "severity": "MAJOR",
                "check": "empty_slide",
                "message": "Slide has no visible text content at all",
            })
        elif total_chars > 2000:
            issues.append({
                "slide": slide_idx,
                "severity": "MINOR",
                "check": "text_density_high",
                "message": (
                    f"Slide has {total_chars} characters of text - "
                    f"consider splitting into multiple slides"
                ),
            })
    return issues


# ── Main runner ───────────────────────────────────────────────────────────────

def run_all_checks(pptx_path: str, expected_slides: int | None = None) -> dict:
    """Run all QA checks and return a structured report."""
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {
            "status": "ERROR",
            "file": pptx_path,
            "error": str(e),
            "issues": [],
            "summary": {"CRITICAL": 0, "MAJOR": 0, "MINOR": 0},
        }

    all_issues = []
    checks = [
        ("slide_count", lambda: check_slide_count(prs, expected_slides)),
        ("shape_overflow", lambda: check_shape_overflow(prs)),
        ("placeholder_text", lambda: check_placeholder_text(prs)),
        ("speaker_notes", lambda: check_speaker_notes(prs)),
        ("font_sizes", lambda: check_font_sizes(prs)),
        ("text_overflow", lambda: check_text_overflow_heuristic(prs)),
        ("shape_overlap", lambda: check_shape_overlap(prs)),
        ("empty_text_frames", lambda: check_empty_text_frames(prs)),
        ("closing_cta", lambda: check_closing_slide_cta_count(prs)),
        ("content_margins", lambda: check_content_margins(prs)),
        ("text_density", lambda: check_slide_text_density(prs)),
    ]

    for check_name, check_fn in checks:
        try:
            issues = check_fn()
            all_issues.extend(issues)
        except Exception as e:
            all_issues.append({
                "slide": 0,
                "severity": "MINOR",
                "check": check_name,
                "message": f"Check failed with error: {e}",
            })

    # Summarize
    summary = defaultdict(int)
    for issue in all_issues:
        summary[issue["severity"]] += 1

    # Group by slide
    by_slide: dict[int, list[dict]] = defaultdict(list)
    for issue in all_issues:
        by_slide[issue["slide"]].append(issue)

    has_critical_or_major = summary["CRITICAL"] > 0 or summary["MAJOR"] > 0
    status = "ISSUES_FOUND" if has_critical_or_major else "CLEAN"

    return {
        "status": status,
        "file": pptx_path,
        "slide_count": len(prs.slides),
        "expected_slides": expected_slides,
        "issues": all_issues,
        "issues_by_slide": {k: v for k, v in sorted(by_slide.items())},
        "summary": dict(summary),
    }


def format_report(report: dict) -> str:
    """Format the report as human-readable text."""
    lines = []
    lines.append(f"## PPTX QA Report")
    lines.append(f"")
    lines.append(f"**Status:** {report['status']}")
    lines.append(f"**File:** {report['file']}")
    lines.append(f"**Slides:** {report.get('slide_count', '?')}")
    if report.get("expected_slides"):
        lines.append(f"**Expected:** {report['expected_slides']}")
    lines.append(f"")
    lines.append(f"### Summary")
    summary = report.get("summary", {})
    lines.append(f"- CRITICAL: {summary.get('CRITICAL', 0)}")
    lines.append(f"- MAJOR: {summary.get('MAJOR', 0)}")
    lines.append(f"- MINOR: {summary.get('MINOR', 0)}")
    lines.append(f"")

    if not report.get("issues"):
        lines.append("No issues found.")
    else:
        lines.append(f"### Issues by Slide")
        lines.append(f"")
        by_slide = report.get("issues_by_slide", {})
        for slide_num in sorted(by_slide.keys()):
            slide_issues = by_slide[slide_num]
            if slide_num == 0:
                lines.append(f"#### General")
            else:
                lines.append(f"#### Slide {slide_num}")
            for issue in slide_issues:
                lines.append(f"- **[{issue['severity']}]** ({issue['check']}) {issue['message']}")
            lines.append(f"")

    return "\n".join(lines)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PPTX QA checks")
    parser.add_argument("pptx_path", help="Path to .pptx file")
    parser.add_argument("--expected-slides", type=int, default=None,
                        help="Expected number of slides")
    parser.add_argument("--json", action="store_true",
                        help="Output JSON instead of text")
    args = parser.parse_args()

    report = run_all_checks(args.pptx_path, args.expected_slides)

    if args.json:
        # Convert defaultdict keys for JSON serialization
        print(json.dumps(report, indent=2, default=str))
    else:
        print(format_report(report))

    sys.exit(0 if report["status"] == "CLEAN" else
             2 if report["status"] == "ERROR" else 1)
