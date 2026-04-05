"""FastAPI server exposing CSA-Copilot over HTTP + WebSocket.

Start with:
    python app.py --server [--port PORT]

The server picks a random free port if ``--port`` is omitted and writes the
chosen port to stdout as ``PORT:XXXX`` so the Electron main process can read it.

Security:
  - Binds only to 127.0.0.1 (never 0.0.0.0).
  - The /file endpoint validates all paths are under ``outputs/``.
  - WebSocket messages are validated before processing.
  - No secrets are exposed through any endpoint.
"""

from __future__ import annotations

import asyncio
import base64
import contextlib
import io
import json
import logging
import os
import time
from pathlib import Path
from typing import Any

from fastapi import FastAPI, HTTPException, WebSocket, WebSocketDisconnect
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# App-wide state injected at startup by server_main()
# ---------------------------------------------------------------------------

_app_dir: Path = Path(__file__).resolve().parent
_outputs_dir: Path = _app_dir / "outputs"
_event_store: Any = None        # EventStore instance
_copilot_client: Any = None     # CopilotClient instance
_session_map: dict[str, Any] = {}   # session_id -> Session object
_collector: Any = None          # EventCollector instance

# ---------------------------------------------------------------------------
# FastAPI application
# ---------------------------------------------------------------------------

app = FastAPI(title="CSA Copilot API", version="1.0.0", docs_url=None, redoc_url=None)

# Allow the Electron renderer (loaded from file://) to call the API.
# Restricted to only localhost origins for security.
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost",
        "http://127.0.0.1",
        "http://localhost:5173",
        "http://localhost:3000",
        "http://127.0.0.1:5173",
    ],
    allow_credentials=False,
    allow_methods=["GET", "POST", "DELETE"],
    allow_headers=["Content-Type"],
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _store() -> Any:
    if _event_store is None:
        raise HTTPException(status_code=503, detail="Store not initialised")
    return _event_store


def _safe_outputs_path(raw: str) -> Path:
    """Resolve *raw* and assert it is inside outputs/.

    Raises HTTPException 400 if the path escapes the outputs directory.
    """
    if not raw or not raw.strip():
        raise HTTPException(status_code=400, detail="Invalid path")
    if "\x00" in raw:
        raise HTTPException(status_code=400, detail="Invalid path")

    outputs_resolved = _outputs_dir.resolve()

    try:
        raw_path = Path(raw)
        # Always resolve relative to the outputs directory, never CWD.
        if raw_path.is_absolute():
            resolved = raw_path.resolve()
        else:
            resolved = (outputs_resolved / raw_path).resolve()
    except (ValueError, OSError):
        raise HTTPException(status_code=400, detail="Invalid path")

    if not resolved.is_relative_to(outputs_resolved):
        raise HTTPException(status_code=400, detail="Path outside outputs directory")
    return resolved


# ---------------------------------------------------------------------------
# Health
# ---------------------------------------------------------------------------

@app.get("/health")
async def health() -> dict[str, str]:
    return {"status": "ok", "version": "1.0.0"}


# ---------------------------------------------------------------------------
# Agents
# ---------------------------------------------------------------------------

@app.get("/agents")
async def list_agents() -> JSONResponse:
    from agents import CATALOG
    result = []
    for name, agent in CATALOG.all_agents.items():
        result.append({
            "name": name,
            "display_name": agent.display_name,
            "description": agent.description,
            "model": getattr(agent, "model", ""),
            "infer": agent.infer,
        })
    return JSONResponse(content=result)


# ---------------------------------------------------------------------------
# Sessions
# ---------------------------------------------------------------------------

class CreateSessionRequest(BaseModel):
    agent: str | None = None
    model: str | None = None


@app.post("/sessions")
async def create_session(body: CreateSessionRequest) -> JSONResponse:
    from agents import ALL_AGENT_CONFIGS, ALL_SKILL_DIRS, DEFAULT_MODEL
    from tools import ALL_CUSTOM_TOOLS

    if _copilot_client is None:
        raise HTTPException(status_code=503, detail="Copilot client not ready")

    model = body.model or DEFAULT_MODEL
    _sid_ref: list[str] = []  # mutable container; filled after create_session

    async def _perm(request: Any, _inv: Any) -> Any:
        from copilot.types import PermissionRequestResult
        return PermissionRequestResult(kind="approved")

    async def _user_input(request: Any, _inv: Any) -> Any:
        from copilot.types import UserInputResponse
        from server_adapter import pop_user_response, _send, set_pending_input
        question = request.get("question", "")
        choices = request.get("choices")
        sid = _sid_ref[0] if _sid_ref else None
        payload = {"type": "waiting_for_input", "question": question, "choices": choices}
        if sid:
            set_pending_input(sid, payload)
        _send(payload, sid)
        try:
            answer = await pop_user_response(timeout=300.0, session_id=sid)
        except asyncio.TimeoutError:
            answer = ""
        finally:
            if sid:
                set_pending_input(sid, None)
                _send({"type": "input_resolved"}, sid)
        return UserInputResponse(answer=answer, wasFreeform=True)

    try:
        session = await _copilot_client.create_session(
            {
                "model": model,
                "streaming": True,
                "custom_agents": ALL_AGENT_CONFIGS,
                "tools": ALL_CUSTOM_TOOLS,
                "skill_directories": ALL_SKILL_DIRS,
                "on_permission_request": _perm,
                "on_user_input_request": _user_input,
                "working_directory": str(_app_dir),
            }
        )
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc))

    sid = session.session_id
    _sid_ref.append(sid)  # now the _user_input callback can find it
    _session_map[sid] = session

    if _collector:
        _collector.on_session_created(
            sid, agent=body.agent or "", model=model, frontend="desktop"
        )
    if _event_store and body.agent:
        _event_store.update_session_agent(sid, body.agent)

    return JSONResponse(content={"session_id": sid, "model": model})


class ResumeSessionRequest(BaseModel):
    pass


@app.post("/sessions/{session_id}/resume")
async def resume_session(session_id: str, body: ResumeSessionRequest) -> JSONResponse:
    from agents import ALL_AGENT_CONFIGS, ALL_SKILL_DIRS
    from tools import ALL_CUSTOM_TOOLS

    if _copilot_client is None:
        raise HTTPException(status_code=503, detail="Copilot client not ready")
    if _event_store is None:
        raise HTTPException(status_code=503, detail="Store not initialised")

    full_id = _event_store.resolve_prefix("sessions", session_id)
    if not full_id:
        raise HTTPException(status_code=404, detail="Session not found")

    s_detail = _event_store.get_session(full_id)
    if not s_detail or not s_detail.get("resumable"):
        raise HTTPException(status_code=400, detail="Session is not resumable")

    async def _perm(request: Any, _inv: Any) -> Any:
        from copilot.types import PermissionRequestResult
        return PermissionRequestResult(kind="approved")

    async def _user_input(request: Any, _inv: Any) -> Any:
        from copilot.types import UserInputResponse
        from server_adapter import _send, pop_user_response, set_pending_input
        question = request.get("question", "")
        choices = request.get("choices")
        payload = {"type": "waiting_for_input", "question": question, "choices": choices}
        set_pending_input(full_id, payload)
        _send(payload, full_id)
        try:
            answer = await pop_user_response(timeout=300.0, session_id=full_id)
        except asyncio.TimeoutError:
            answer = ""
        finally:
            set_pending_input(full_id, None)
            _send({"type": "input_resolved"}, full_id)
        return UserInputResponse(answer=answer, wasFreeform=True)

    try:
        session = await _copilot_client.resume_session(
            full_id,
            {
                "streaming": True,
                "custom_agents": ALL_AGENT_CONFIGS,
                "tools": ALL_CUSTOM_TOOLS,
                "skill_directories": ALL_SKILL_DIRS,
                "on_permission_request": _perm,
                "on_user_input_request": _user_input,
                "working_directory": str(_app_dir),
            },
        )
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc))

    _session_map[full_id] = session
    _event_store.reactivate_session(full_id)

    s_detail = _event_store.get_session(full_id)
    turns = _event_store.get_turns(full_id)

    return JSONResponse(content={
        "session_id": full_id,
        "agent": s_detail.get("agent", ""),
        "model": s_detail.get("model", ""),
        "turn_count": s_detail.get("turn_count", 0),
        "turns": turns,
    })


@app.delete("/sessions/{session_id}")
async def end_session(session_id: str) -> JSONResponse:
    from server_adapter import unregister_event_handler
    if _event_store:
        full_id = _event_store.resolve_prefix("sessions", session_id) or session_id
        _event_store.end_session(full_id, resumable=False)
    session = _session_map.pop(session_id, None)
    if session:
        unregister_event_handler(session_id)
        with contextlib.suppress(Exception):
            session._event_handlers.clear()
    return JSONResponse(content={"ok": True})


@app.get("/sessions")
async def list_sessions(all: bool = False) -> JSONResponse:
    store = _store()
    if all:
        sessions = store.list_sessions()
    else:
        sessions = store.list_sessions(status="active")
    return JSONResponse(content=sessions)


@app.get("/sessions/{session_id}")
async def get_session(session_id: str) -> JSONResponse:
    store = _store()
    full_id = store.resolve_prefix("sessions", session_id) or session_id
    detail = store.get_session(full_id)
    if not detail:
        raise HTTPException(status_code=404, detail="Session not found")
    return JSONResponse(content=detail)


@app.get("/sessions/{session_id}/turns")
async def get_session_turns(session_id: str) -> JSONResponse:
    store = _store()
    full_id = store.resolve_prefix("sessions", session_id) or session_id
    turns = store.get_turns(full_id)
    return JSONResponse(content=turns)


@app.get("/sessions/{session_id}/turns/{turn_number}/invocations")
async def get_turn_invocations(session_id: str, turn_number: int) -> JSONResponse:
    store = _store()
    full_id = store.resolve_prefix("sessions", session_id) or session_id
    turns = store.get_turns(full_id)
    turn = next((t for t in turns if t.get("turn_number") == turn_number), None)
    if not turn:
        raise HTTPException(status_code=404, detail="Turn not found")
    invocations = store.get_invocations(turn["id"])
    return JSONResponse(content=invocations)


@app.get("/sessions/{session_id}/status")
async def get_session_status(session_id: str) -> JSONResponse:
    """Lightweight status check for frontend reconnect polling."""
    from server_adapter import _ws_map

    store = _store()
    full_id = store.resolve_prefix("sessions", session_id)
    if not full_id:
        raise HTTPException(status_code=404, detail="Session not found")

    detail = store.get_session(full_id)
    if not detail:
        raise HTTPException(status_code=404, detail="Session not found")

    in_memory = full_id in _session_map
    has_ws = bool(_ws_map.get(full_id))

    return JSONResponse(content={
        "session_id": full_id,
        "status": detail.get("status", "ended"),
        "in_memory": in_memory,
        "has_running_turn": in_memory and has_ws,
        "turn_count": detail.get("turn_count", 0),
        "resumable": bool(detail.get("resumable", 0)),
    })


@app.get("/sessions/{session_id}/events")
async def get_session_events(session_id: str) -> JSONResponse:
    """Historical events reconstructed from DB for the frontend activity feed."""
    store = _store()
    full_id = store.resolve_prefix("sessions", session_id)
    if not full_id:
        raise HTTPException(status_code=404, detail="Session not found")

    detail = store.get_session(full_id)
    if not detail:
        raise HTTPException(status_code=404, detail="Session not found")

    turns = store.get_turns(full_id)
    events: list[dict[str, Any]] = []

    for i, turn in enumerate(turns):
        # turn_started
        events.append({
            "type": "turn_started",
            "data": {
                "turn_number": turn.get("turn_number", i + 1),
                "agent": turn.get("agent", ""),
                "user_prompt": turn.get("user_prompt", ""),
            },
            "time": turn.get("started_at", ""),
        })

        # invocations
        invocations = store.get_invocations_for_turn(turn["id"])
        for inv in invocations:
            inv_type = inv.get("type", "")
            inv_name = inv.get("name", "")

            if inv_type == "tool_call":
                events.append({
                    "type": "tool_started",
                    "data": {"tool": inv_name, "args": inv.get("input", "{}")},
                    "time": inv.get("started_at", ""),
                })
                output_raw = inv.get("output") or ""
                events.append({
                    "type": "tool_completed",
                    "data": {
                        "tool": inv_name,
                        "output_preview": output_raw[:500],
                        "duration_ms": inv.get("duration_ms", 0),
                    },
                    "time": inv.get("ended_at") or inv.get("started_at", ""),
                })
            elif inv_type == "subagent":
                events.append({
                    "type": "subagent_started",
                    "data": {"agent": inv_name},
                    "time": inv.get("started_at", ""),
                })
                events.append({
                    "type": "subagent_completed",
                    "data": {
                        "agent": inv_name,
                        "duration_ms": inv.get("duration_ms", 0),
                    },
                    "time": inv.get("ended_at") or inv.get("started_at", ""),
                })

        # usage
        input_tok = turn.get("input_tokens", 0)
        output_tok = turn.get("output_tokens", 0)
        if input_tok or output_tok:
            events.append({
                "type": "usage",
                "data": {
                    "turn_number": turn.get("turn_number", i + 1),
                    "input_tokens": input_tok,
                    "output_tokens": output_tok,
                    "cache_read_tokens": turn.get("cache_read_tokens", 0),
                    "cache_write_tokens": turn.get("cache_write_tokens", 0),
                    "estimated_cost_usd": turn.get("estimated_cost_usd", 0.0),
                },
                "time": turn.get("ended_at") or turn.get("started_at", ""),
            })

    # done event for ended sessions
    if detail.get("status") == "ended":
        events.append({
            "type": "done",
            "data": {"session_id": full_id, "reason": "session_ended"},
            "time": detail.get("ended_at") or detail.get("started_at", ""),
        })

    return JSONResponse(content=events)


# ---------------------------------------------------------------------------
# Usage
# ---------------------------------------------------------------------------

@app.get("/usage")
async def get_usage(
    period: str = "all",
    agent: str | None = None,
    model: str | None = None,
) -> JSONResponse:
    from queries import usage_summary
    store = _store()
    data = usage_summary(store, period=period if period != "all" else None, agent=agent, model=model)
    return JSONResponse(content=data)


# ---------------------------------------------------------------------------
# Outputs
# ---------------------------------------------------------------------------

_FILE_TYPE_MAP = {
    ".pptx": "pptx",
    ".md": "markdown",
    ".py": "python",
    ".sh": "shell",
    ".bicep": "bicep",
    ".json": "json",
    ".yaml": "yaml",
    ".yml": "yaml",
    ".ts": "typescript",
    ".txt": "text",
}

_SKIP_DIRS = {".fragments", "__pycache__", ".git", "node_modules"}


def _classify_output_category(path: Path) -> str:
    parts = path.parts
    for i, part in enumerate(parts):
        if part == "outputs" and i + 1 < len(parts):
            return parts[i + 1]
    return "other"


@app.get("/outputs")
async def list_outputs() -> JSONResponse:
    outputs_resolved = _outputs_dir.resolve()
    results = []
    for p in outputs_resolved.rglob("*"):
        if not p.is_file():
            continue
        if any(part in _SKIP_DIRS for part in p.parts):
            continue
        suffix = p.suffix.lower()
        if suffix not in _FILE_TYPE_MAP:
            continue
        if "-plan.md" in p.name:
            continue
        try:
            stat = p.stat()
        except OSError:
            continue
        rel = str(p.relative_to(outputs_resolved))
        results.append({
            "path": str(p),
            "relative": rel,
            "name": p.name,
            "type": _FILE_TYPE_MAP.get(suffix, "file"),
            "category": _classify_output_category(p),
            "size": stat.st_size,
            "modified": stat.st_mtime,
        })
    results.sort(key=lambda x: x["modified"], reverse=True)
    return JSONResponse(content=results)


@app.get("/outputs/grouped")
async def list_outputs_grouped() -> JSONResponse:
    """Return outputs grouped as logical deliverables.

    - slides: each .pptx is a deliverable (with companion .pdf, generate_*.py)
    - demos: each *-demos.md + its companion folder is a deliverable
    - hackathons: each subfolder is a deliverable
    - ai-projects: each subfolder is a deliverable
    """
    import re as _re
    outputs_resolved = _outputs_dir.resolve()
    groups: list[dict] = []

    # ── Slides ────────────────────────────────────────────────
    slides_dir = outputs_resolved / "slides"
    if slides_dir.is_dir():
        seen_stems: set[str] = set()
        try:
            slide_entries = sorted(slides_dir.iterdir())
        except (PermissionError, OSError):
            slide_entries = []
        for p in slide_entries:
            if not p.is_file() or p.suffix.lower() != ".pptx":
                continue
            if any(part in _SKIP_DIRS for part in p.parts):
                continue
            stem = p.stem
            if stem in seen_stems:
                continue
            seen_stems.add(stem)

            # Companion files
            companions = [str(p)]
            pdf = slides_dir / f"{stem}.pdf"
            if pdf.is_file():
                companions.append(str(pdf))
            gen = slides_dir / f"generate_{stem.replace('-', '_')}_pptx.py"
            if gen.is_file():
                companions.append(str(gen))

            # Parse metadata from filename
            level_m = _re.search(r"[_-](l[1-4]00)[_-]", stem, _re.IGNORECASE)
            dur_m = _re.search(r"(\d+)\s*(?:min|h)", stem, _re.IGNORECASE)
            title = stem.replace("-", " ").replace("_", " ")
            title = _re.sub(r"\bl\d{3}\b", "", title, flags=_re.IGNORECASE).strip()
            title = _re.sub(r"\d+\s*(?:min|h)\b", "", title, flags=_re.IGNORECASE).strip()
            title = " ".join(w.capitalize() for w in title.split())

            try:
                stat = p.stat()
            except OSError:
                continue

            groups.append({
                "id": f"slides/{stem}",
                "title": title or stem,
                "category": "slides",
                "primary_file": str(p),
                "file_count": len(companions),
                "files": companions,
                "content_level": level_m.group(1).upper() if level_m else None,
                "duration": dur_m.group(0) if dur_m else None,
                "has_pdf": pdf.is_file(),
                "size": stat.st_size,
                "modified": stat.st_mtime,
            })

    # ── Demos ─────────────────────────────────────────────────
    demos_dir = outputs_resolved / "demos"
    if demos_dir.is_dir():
        try:
            demo_entries = sorted(demos_dir.iterdir())
        except (PermissionError, OSError):
            demo_entries = []
        for p in demo_entries:
            if not p.is_file() or not p.name.endswith("-demos.md"):
                continue
            slug = p.name.replace("-demos.md", "")
            companion_dir = demos_dir / slug
            file_list = [str(p)]
            if companion_dir.is_dir():
                for child in companion_dir.rglob("*"):
                    if child.is_file() and not any(part in _SKIP_DIRS for part in child.parts):
                        file_list.append(str(child))

            title = slug.replace("-", " ").replace("generic ", "").replace("internal ", "")
            title = " ".join(w.capitalize() for w in title.split())

            try:
                stat = p.stat()
            except OSError:
                continue

            groups.append({
                "id": f"demos/{slug}",
                "title": title or slug,
                "category": "demos",
                "primary_file": str(p),
                "file_count": len(file_list),
                "files": file_list,
                "content_level": None,
                "duration": None,
                "has_pdf": False,
                "size": sum(Path(f).stat().st_size for f in file_list if Path(f).is_file()),
                "modified": stat.st_mtime,
            })

    # ── Hackathons ────────────────────────────────────────────
    hack_dir = outputs_resolved / "hackathons"
    if hack_dir.is_dir():
        try:
            hack_entries = sorted(hack_dir.iterdir())
        except (PermissionError, OSError):
            hack_entries = []
        for d in hack_entries:
            if not d.is_dir() or d.name.startswith("."):
                continue
            file_list = []
            for child in d.rglob("*"):
                if child.is_file() and not any(part in _SKIP_DIRS for part in child.parts):
                    file_list.append(str(child))
            if not file_list:
                continue

            title = d.name.replace("-", " ")
            title = " ".join(w.capitalize() for w in title.split())
            readme = d / "README.md"
            latest = max(Path(f).stat().st_mtime for f in file_list if Path(f).is_file())

            groups.append({
                "id": f"hackathons/{d.name}",
                "title": title,
                "category": "hackathons",
                "primary_file": str(readme) if readme.is_file() else file_list[0],
                "file_count": len(file_list),
                "files": file_list,
                "content_level": None,
                "duration": None,
                "has_pdf": False,
                "size": sum(Path(f).stat().st_size for f in file_list if Path(f).is_file()),
                "modified": latest,
            })

    # ── AI Projects ───────────────────────────────────────────
    _ARCH_DOC_NAMES = {"solution-design.md", "architecture-diagram.md", "architecture.md"}
    proj_dir = outputs_resolved / "ai-projects"
    if proj_dir.is_dir():
        try:
            proj_entries = sorted(proj_dir.iterdir())
        except (PermissionError, OSError):
            proj_entries = []
        for d in proj_entries:
            if not d.is_dir() or d.name.startswith("."):
                continue
            file_list = []
            for child in d.rglob("*"):
                if child.is_file() and not any(
                    part in _SKIP_DIRS | {"__pycache__", ".pytest_cache"}
                    for part in child.parts
                ):
                    file_list.append(str(child))
            if not file_list:
                continue

            title = d.name.replace("-", " ").replace("_", " ")
            title = " ".join(w.capitalize() for w in title.split())
            readme = d / "README.md"
            latest = max(Path(f).stat().st_mtime for f in file_list if Path(f).is_file())

            # Check for architecture documents in docs/
            docs_dir = d / "docs"
            arch_docs: list[str] = []
            if docs_dir.is_dir():
                for ad in docs_dir.iterdir():
                    if ad.is_file() and ad.name.lower() in _ARCH_DOC_NAMES:
                        arch_docs.append(str(ad))

            groups.append({
                "id": f"ai-projects/{d.name}",
                "title": title,
                "category": "ai-projects",
                "primary_file": str(readme) if readme.is_file() else file_list[0],
                "file_count": len(file_list),
                "files": file_list,
                "content_level": None,
                "duration": None,
                "has_pdf": False,
                "has_architecture": len(arch_docs) > 0,
                "architecture_docs": arch_docs,
                "size": sum(Path(f).stat().st_size for f in file_list if Path(f).is_file()),
                "modified": latest,
            })

    groups.sort(key=lambda x: x["modified"], reverse=True)
    return JSONResponse(content=groups)


# ---------------------------------------------------------------------------
# Delete output
# ---------------------------------------------------------------------------

@app.delete("/outputs")
async def delete_output(path: str) -> JSONResponse:
    """Delete a file or directory under outputs/."""
    import shutil
    resolved = _safe_outputs_path(path)
    if resolved.is_file():
        resolved.unlink()
    elif resolved.is_dir():
        shutil.rmtree(str(resolved))
    else:
        raise HTTPException(status_code=404, detail="Not found")
    # Return only the relative path — never expose absolute server paths.
    rel = str(resolved.relative_to(_outputs_dir.resolve()))
    return JSONResponse(content={"ok": True, "deleted": rel})


@app.delete("/outputs/grouped")
async def delete_grouped_output(id: str) -> JSONResponse:
    """Delete all files belonging to a grouped output.

    The *id* follows the pattern ``category/slug`` (e.g. ``slides/my-deck``,
    ``hackathons/azure-ai``, ``ai-projects/contoso``).  For directory-based
    groups (hackathons, ai-projects) we remove the entire subfolder.  For
    slides and demos we remove the individual companion files.
    """
    import shutil

    parts = id.strip("/").split("/", 1)
    if len(parts) != 2 or not parts[0] or not parts[1]:
        raise HTTPException(status_code=400, detail="Invalid group id")

    category, slug = parts
    outputs_resolved = _outputs_dir.resolve()
    deleted: list[str] = []

    if category in ("hackathons", "ai-projects"):
        # Directory-based groups — remove the whole subfolder
        target = outputs_resolved / category / slug
        resolved = target.resolve()
        try:
            resolved.relative_to(outputs_resolved)
        except ValueError:
            raise HTTPException(status_code=400, detail="Path outside outputs directory")
        if resolved.is_dir():
            shutil.rmtree(str(resolved))
            deleted.append(str(resolved))
        else:
            raise HTTPException(status_code=404, detail="Directory not found")

    elif category == "slides":
        # Slide groups: pptx + optional pdf + optional generate_*.py
        slides_dir = outputs_resolved / "slides"
        for p in slides_dir.iterdir():
            if p.is_file() and (
                p.stem == slug
                or p.stem == f"{slug}"
                or p.name == f"generate_{slug.replace('-', '_')}_pptx.py"
            ):
                resolved = p.resolve()
                try:
                    resolved.relative_to(outputs_resolved)
                except ValueError:
                    continue
                resolved.unlink()
                deleted.append(str(resolved))

    elif category == "demos":
        # Demo groups: {slug}-demos.md + optional companion directory
        demos_dir = outputs_resolved / "demos"
        md_file = demos_dir / f"{slug}-demos.md"
        if md_file.is_file():
            md_file.unlink()
            deleted.append(str(md_file))
        companion = demos_dir / slug
        if companion.is_dir():
            shutil.rmtree(str(companion))
            deleted.append(str(companion))
        if not deleted:
            raise HTTPException(status_code=404, detail="Demo files not found")

    else:
        raise HTTPException(status_code=400, detail=f"Unknown category: {category}")

    return JSONResponse(content={"ok": True, "deleted": deleted})


@app.get("/file")
async def read_file(path: str) -> JSONResponse:
    resolved = _safe_outputs_path(path)
    if not resolved.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    try:
        content = resolved.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc))
    # Return only the relative path — never expose absolute server paths.
    rel = str(resolved.relative_to(_outputs_dir.resolve()))
    return JSONResponse(content={"path": rel, "content": content})


# ---------------------------------------------------------------------------
# File download (raw binary with Content-Disposition)
# ---------------------------------------------------------------------------

@app.get("/file/download")
async def download_file(path: str):
    from fastapi.responses import FileResponse
    resolved = _safe_outputs_path(path)
    if not resolved.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(
        path=str(resolved),
        filename=resolved.name,
        media_type="application/octet-stream",
    )


# ---------------------------------------------------------------------------
# ZIP export
# ---------------------------------------------------------------------------

class ZipRequest(BaseModel):
    paths: list[str]
    name: str | None = None


@app.post("/outputs/zip")
async def create_zip(body: ZipRequest):
    import zipfile
    from fastapi.responses import StreamingResponse

    if not body.paths:
        raise HTTPException(status_code=400, detail="No paths provided")

    # Validate all paths
    resolved_paths: list[Path] = []
    for p in body.paths:
        resolved_paths.append(_safe_outputs_path(p))

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for rp in resolved_paths:
            if rp.is_file():
                zf.write(str(rp), rp.name)
            elif rp.is_dir():
                for child in rp.rglob("*"):
                    if child.is_file():
                        arcname = str(child.relative_to(rp.parent))
                        zf.write(str(child), arcname)

    buf.seek(0)
    # Sanitise user-supplied name to prevent header injection.
    import re as _re_zip
    safe_name = _re_zip.sub(r'[^\w\-.]', '_', body.name or "csa-copilot-export")
    zip_name = safe_name + ".zip"
    return StreamingResponse(
        buf,
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="{zip_name}"'},
    )


# ---------------------------------------------------------------------------
# Output metadata
# ---------------------------------------------------------------------------

@app.get("/outputs/metadata")
async def get_output_metadata(path: str) -> JSONResponse:
    """Parse structured metadata from an output file or its companion plan."""
    import re as _re
    resolved = _safe_outputs_path(path)
    if not resolved.exists():
        raise HTTPException(status_code=404, detail="Not found")

    name = resolved.name
    meta: dict[str, Any] = {
        "title": name.rsplit(".", 1)[0] if "." in name else name,
        "category": _classify_output_category(resolved),
        "size": resolved.stat().st_size if resolved.is_file() else 0,
        "modified": resolved.stat().st_mtime if resolved.is_file() else 0,
    }

    # Parse level from filename  e.g.  "keda-banking-l300-30min.pptx"
    m = _re.search(r"[_-](l[1-4]00)[_-]", name, _re.IGNORECASE)
    if m:
        meta["contentLevel"] = m.group(1).upper()

    # Parse duration from filename
    m = _re.search(r"(\d+)\s*min", name, _re.IGNORECASE)
    if m:
        meta["duration"] = f"{m.group(1)} min"
    else:
        m = _re.search(r"(\d+)\s*h", name, _re.IGNORECASE)
        if m:
            meta["duration"] = f"{m.group(1)}h"

    # Slide count for pptx files
    if resolved.suffix.lower() == ".pptx" and resolved.is_file():
        try:
            from pptx import Presentation as _Prs
            prs = _Prs(str(resolved))
            meta["slideCount"] = len(prs.slides)
        except Exception:
            pass

    # Try to find companion plan file
    plans_dir = _app_dir / "plans"
    if plans_dir.is_dir():
        stem = name.rsplit(".", 1)[0]
        for suffix in ["-complete.md", "-plan.md"]:
            plan = plans_dir / (stem + suffix)
            if plan.is_file():
                # Return only the relative path — never expose absolute server paths.
                meta["planFile"] = f"plans/{plan.name}"
                break

    return JSONResponse(content=meta)


# ---------------------------------------------------------------------------
# PPTX Preview — renders actual slide content via LibreOffice + PyMuPDF
# ---------------------------------------------------------------------------

class PptxPreviewRequest(BaseModel):
    path: str
    max_width: int = 1280


@app.post("/preview/pptx")
async def preview_pptx(body: PptxPreviewRequest) -> JSONResponse:
    """Convert PPTX → PDF (LibreOffice) → PNG per page (PyMuPDF).

    Caches rendered PNGs in outputs/slides/.img_cache/{hash}/.
    Falls back to text-only placeholders if LibreOffice is unavailable.
    """
    import hashlib
    import shutil
    import subprocess
    import tempfile

    resolved = _safe_outputs_path(body.path)
    if not resolved.is_file() or resolved.suffix.lower() != ".pptx":
        raise HTTPException(status_code=400, detail="Not a .pptx file")

    # ── Cache key ─────────────────────────────────────────────────────────
    file_hash = hashlib.md5(
        f"{resolved.name}:{resolved.stat().st_mtime}:{body.max_width}".encode()
    ).hexdigest()[:12]
    cache_dir = _outputs_dir / "slides" / ".img_cache" / file_hash
    cache_dir.mkdir(parents=True, exist_ok=True)

    # ── Check cache ───────────────────────────────────────────────────────
    cached_pngs = sorted(cache_dir.glob("slide-*.png"))

    if not cached_pngs:
        # Convert PPTX → PDF → per-page PNGs
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            raise HTTPException(status_code=500, detail="LibreOffice not installed")

        with tempfile.TemporaryDirectory() as tmpdir:
            # Step 1: PPTX → PDF
            try:
                subprocess.run(
                    [soffice, "--headless", "--norestore", "--convert-to", "pdf",
                     "--outdir", tmpdir, str(resolved)],
                    capture_output=True, timeout=120, check=True,
                    env={**os.environ, "HOME": tmpdir},  # avoid lock conflicts
                )
            except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as exc:
                raise HTTPException(status_code=500, detail=f"LibreOffice conversion failed: {exc}")

            pdf_files = list(Path(tmpdir).glob("*.pdf"))
            if not pdf_files:
                raise HTTPException(status_code=500, detail="PDF conversion produced no output")

            # Step 2: PDF → PNG per page using PyMuPDF
            try:
                import fitz  # PyMuPDF
            except ImportError:
                raise HTTPException(status_code=500, detail="PyMuPDF not installed")

            pdf_doc = fitz.open(str(pdf_files[0]))
            for page_num in range(len(pdf_doc)):
                page = pdf_doc[page_num]
                # Scale to requested width
                scale = body.max_width / page.rect.width
                mat = fitz.Matrix(scale, scale)
                pix = page.get_pixmap(matrix=mat)
                out_path = cache_dir / f"slide-{page_num:03d}.png"
                pix.save(str(out_path))
            pdf_doc.close()

        cached_pngs = sorted(cache_dir.glob("slide-*.png"))

    if not cached_pngs:
        raise HTTPException(status_code=500, detail="No slides rendered")

    # ── Extract notes from PPTX ───────────────────────────────────────────
    notes_list: list[str] = []
    titles_list: list[str] = []
    try:
        from pptx import Presentation
        prs = Presentation(str(resolved))
        for slide in list(prs.slides):
            notes_text = ""
            title_text = ""
            try:
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
            except Exception:
                pass
            # Extract title
            texts_found: list[str] = []
            try:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    text = shape.text_frame.text.strip()
                    if not text:
                        continue
                    is_title = False
                    try:
                        ph = shape.placeholder_format
                        if ph and ph.idx in (0, 15):
                            is_title = True
                    except (ValueError, AttributeError):
                        pass
                    if not is_title and hasattr(shape, "name") and "title" in shape.name.lower():
                        is_title = True
                    if is_title and not title_text:
                        title_text = text
                    else:
                        texts_found.append(text[:200])
                if not title_text and texts_found:
                    candidates = texts_found[:3]
                    title_text = max(candidates, key=len)
            except Exception:
                pass
            notes_list.append(notes_text[:2000])
            titles_list.append(title_text)
    except Exception:
        pass

    # ── Build response ────────────────────────────────────────────────────
    slides_data = []
    for i, png_path in enumerate(cached_pngs):
        png_b64 = base64.b64encode(png_path.read_bytes()).decode("ascii")
        # Get dimensions from the PNG
        try:
            from PIL import Image as _PILImage
            with _PILImage.open(png_path) as img:
                w, h = img.size
        except Exception:
            w, h = body.max_width, int(body.max_width * 9 / 16)
        slides_data.append({
            "index": i,
            "title": titles_list[i] if i < len(titles_list) else "",
            "body_preview": "",
            "notes": notes_list[i] if i < len(notes_list) else "",
            "png_base64": png_b64,
            "width": w,
            "height": h,
        })

    return JSONResponse(content={"slides": slides_data, "total": len(slides_data)})


# Hifi endpoint now just redirects to the main one (same pipeline)
@app.post("/preview/pptx-hifi")
async def preview_pptx_hifi(body: PptxPreviewRequest) -> JSONResponse:
    return await preview_pptx(body)


# ---------------------------------------------------------------------------
# WebSocket — agent streaming
# ---------------------------------------------------------------------------

@app.websocket("/ws/{session_id}")
async def ws_agent(websocket: WebSocket, session_id: str) -> None:
    """Bidirectional streaming channel for a single agent session."""
    from server_adapter import (
        add_ws,
        remove_ws,
        push_user_response,
        pop_user_response,
        set_active_ws,
        set_cancel_flag,
        _ws_map,
        ws_reset,
        _send,
        register_event_handler,
        unregister_event_handler,
    )
    from agents import DEFAULT_TIMEOUT
    from router import route_to_agent

    await websocket.accept()
    is_first_ws = session_id not in _ws_map or len(_ws_map.get(session_id, set())) == 0
    add_ws(session_id, websocket)
    if is_first_ws:
        ws_reset(session_id)

    # Look up or create session
    session = _session_map.get(session_id)
    if session is None:
        # Session may exist in DB but not in memory (server restarted).
        # Create a fresh Copilot session to handle new messages.
        if _copilot_client is not None:
            from agents import ALL_AGENT_CONFIGS, ALL_SKILL_DIRS, DEFAULT_MODEL
            from tools import ALL_CUSTOM_TOOLS

            async def _perm(request: Any, _inv: Any) -> Any:
                from copilot.types import PermissionRequestResult
                return PermissionRequestResult(kind="approved")

            async def _user_input(request: Any, _inv: Any) -> Any:
                from copilot.types import UserInputResponse
                from server_adapter import set_pending_input
                question = request.get("question", "")
                choices = request.get("choices")
                payload = {"type": "waiting_for_input", "question": question, "choices": choices}
                set_pending_input(session_id, payload)
                _send(payload, session_id)
                try:
                    answer = await pop_user_response(timeout=300.0, session_id=session_id)
                except asyncio.TimeoutError:
                    answer = ""
                finally:
                    set_pending_input(session_id, None)
                    _send({"type": "input_resolved"}, session_id)
                return UserInputResponse(answer=answer, wasFreeform=True)

            try:
                session = await _copilot_client.create_session({
                    "model": DEFAULT_MODEL,
                    "streaming": True,
                    "custom_agents": ALL_AGENT_CONFIGS,
                    "tools": ALL_CUSTOM_TOOLS,
                    "skill_directories": ALL_SKILL_DIRS,
                    "on_permission_request": _perm,
                    "on_user_input_request": _user_input,
                    "working_directory": str(_app_dir),
                })
                _session_map[session_id] = session
            except Exception as exc:
                await websocket.send_text(json.dumps(
                    {"type": "error", "message": f"Failed to create session: {exc}"}
                ))
                await websocket.close()
                set_active_ws(session_id, None)
                return
        else:
            await websocket.send_text(json.dumps(
                {"type": "error", "message": "Copilot client not ready"}
            ))
            await websocket.close()
            set_active_ws(session_id, None)
            return

    # Register per-session WS event handler (idempotent — safe on reconnect).
    # In the CLI, session.on(handle_event) is called once and persists.
    # Here we match that: one handler per session, registered on first WS connect,
    # surviving across WS reconnections.
    register_event_handler(session_id, session)

    # Re-send pending waiting_for_input state to reconnecting clients
    from server_adapter import get_pending_input
    pending = get_pending_input(session_id)
    if pending:
        try:
            await websocket.send_text(json.dumps(pending))
        except Exception:
            pass

    current_turn_task: asyncio.Task[Any] | None = None

    async def _run_turn(clean: str, agent_name: str | None, turn_id: str | None) -> None:
        nonlocal current_turn_task

        before_time = time.time()
        turn_status = "success"
        try:
            await session.send_and_wait(
                {"prompt": clean}, timeout=DEFAULT_TIMEOUT
            )
        except asyncio.CancelledError:
            turn_status = "cancelled"
        except TimeoutError:
            turn_status = "timeout"
            # Match CLI behavior: timeout is non-fatal. The agent may still be
            # running server-side. The user can keep chatting.
            _send({
                "type": "error",
                "message": (
                    f"Timeout after {DEFAULT_TIMEOUT // 60} min — the agent is "
                    "still running on the server. You can keep chatting; "
                    "it may deliver results on the next turn."
                ),
            }, session_id)
        except Exception as exc:
            turn_status = "error"
            _send({"type": "error", "message": str(exc)}, session_id)
        finally:
            new_files = _find_new_outputs(before_time)
            if new_files:
                _send({
                    "type": "new_files",
                    "files": [str(f) for f in new_files],
                }, session_id)

            _send({"type": "done", "status": turn_status}, session_id)

            if _collector and turn_id:
                _collector.on_turn_end(
                    turn_id,
                    assistant_response="",
                    status=turn_status,
                )

            if current_turn_task is asyncio.current_task():
                current_turn_task = None

    try:
        while True:
            raw = await websocket.receive_text()
            try:
                msg = json.loads(raw)
            except json.JSONDecodeError:
                await websocket.send_text(json.dumps(
                    {"type": "error", "message": "Invalid JSON"}
                ))
                continue

            msg_type = str(msg.get("type", ""))

            if msg_type == "message":
                content = str(msg.get("content", "")).strip()
                if not content:
                    continue

                if current_turn_task and not current_turn_task.done():
                    _send(
                        {
                            "type": "error",
                            "message": "A turn is already running. Respond to the active prompt or cancel it first.",
                        },
                        session_id,
                    )
                    continue

                ws_reset(session_id)

                # Determine agent via router
                try:
                    agent_name = await route_to_agent(session, content)
                except Exception:
                    agent_name = None

                from agents import DEFAULT_MODEL as _dm

                if agent_name:
                    if _event_store:
                        _event_store.update_session_agent(session_id, agent_name)
                        _event_store.update_session_model(session_id, _dm)

                # Strip @mention prefix
                clean = content
                if content.startswith("@") and " " in content:
                    clean = content.split(" ", 1)[1]

                turn_id = None
                if _collector:
                    turn_id = _collector.on_turn_start(
                        session_id,
                        agent=agent_name or "copilot",
                        model=_dm,
                        user_prompt=clean,
                    )

                _send({"type": "turn_started", "agent": agent_name}, session_id)
                current_turn_task = asyncio.create_task(_run_turn(clean, agent_name, turn_id))

            elif msg_type == "cancel":
                set_cancel_flag(True, session_id)
                if current_turn_task and not current_turn_task.done():
                    current_turn_task.cancel()

            elif msg_type == "user_response":
                content = str(msg.get("content", ""))
                push_user_response(content, session_id)

            elif msg_type == "ping":
                _send({"type": "pong"}, session_id)

            else:
                await websocket.send_text(json.dumps(
                    {"type": "error", "message": f"Unknown message type: {msg_type}"}
                ))

    except WebSocketDisconnect:
        pass
    except Exception as exc:
        log.error("WebSocket error: %s", exc)
    finally:
        remove_ws(session_id, websocket)
        # Unregister the event handler only when ALL WebSockets for the session
        # are gone. This matches CLI behavior where the handler persists for
        # the session lifetime.
        remaining = _ws_map.get(session_id)
        if not remaining:
            unregister_event_handler(session_id)


# ---------------------------------------------------------------------------
# Output file detection (mirrors app.py logic)
# ---------------------------------------------------------------------------

_INTERESTING_SUFFIXES = {".pptx", ".md", ".py", ".bicep", ".json", ".yaml", ".sh"}


def _find_new_outputs(since: float) -> list[Path]:
    found: list[Path] = []
    grace = 3.0
    for p in _outputs_dir.rglob("*"):
        if not p.is_file():
            continue
        if any(part in _SKIP_DIRS for part in p.parts):
            continue
        if p.suffix.lower() not in _INTERESTING_SUFFIXES:
            continue
        if "-plan.md" in p.name:
            continue
        try:
            if p.stat().st_mtime >= since - grace:
                found.append(p)
        except OSError:
            pass
    return found


# ---------------------------------------------------------------------------
# Startup injection (called by server_main in app.py)
# ---------------------------------------------------------------------------

def configure(
    *,
    event_store: Any,
    copilot_client: Any,
    collector: Any,
    app_dir: Path,
    outputs_dir: Path,
) -> None:
    """Inject dependencies into the module-level singletons."""
    global _event_store, _copilot_client, _collector, _app_dir, _outputs_dir
    _event_store = event_store
    _copilot_client = copilot_client
    _collector = collector
    _app_dir = app_dir
    _outputs_dir = outputs_dir
