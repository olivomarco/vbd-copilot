#!/usr/bin/env python3
"""
pptx_utils.py - Shared utilities for generating Microsoft-branded PowerPoint presentations.

Provides reusable building blocks: colors, shapes, slide templates, code blocks,
callouts, tables, metric cards, and more. All presentations share the same
Microsoft brand identity; logo and lead background are loaded from assets/.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from collections import namedtuple
import os
import urllib.request

# ═════════════════════════════════════════════════════════════
# CONSTANTS
# ═════════════════════════════════════════════════════════════

# Microsoft Brand Colors
MS_BLUE = RGBColor(0x00, 0x78, 0xD4)
MS_DARK_BLUE = RGBColor(0x24, 0x3A, 0x5E)
MS_LIGHT_BLUE = RGBColor(0x50, 0xE6, 0xFF)
MS_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MS_LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MS_MID_GRAY = RGBColor(0xD2, 0xD2, 0xD2)
MS_DARK_GRAY = RGBColor(0x3D, 0x3D, 0x3D)
MS_TEXT = RGBColor(0x21, 0x21, 0x21)
MS_TEXT_MUTED = RGBColor(0x61, 0x61, 0x61)
MS_CALLOUT_BG = RGBColor(0xEF, 0xF6, 0xFC)
MS_GREEN = RGBColor(0x10, 0x7C, 0x10)
MS_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
MS_PURPLE = RGBColor(0x88, 0x1C, 0x98)
MS_RED = RGBColor(0xD1, 0x34, 0x38)
MS_BLUE_DARKER = RGBColor(0x00, 0x6C, 0xBE)
MS_CODE_BG = RGBColor(0x24, 0x29, 0x2F)
MS_CODE_TEXT = RGBColor(0xEA, 0xF0, 0xF7)
MS_ACCENT_LIGHT = RGBColor(0xCC, 0xE4, 0xFF)

# Semantic status colors (backgrounds for callout variants)
MS_BLUE_LIGHT_BG = RGBColor(0xDE, 0xEC, 0xF9)   # Info / highlight bg
MS_NAVY_LIGHT = RGBColor(0x2E, 0x4A, 0x6E)       # Subtle dark layering
MS_SUCCESS_BG = RGBColor(0xDF, 0xF6, 0xDD)        # Success / positive bg
MS_WARNING_BG = RGBColor(0xFF, 0xF4, 0xE5)        # Warning bg
MS_ERROR_BG = RGBColor(0xFD, 0xE7, 0xE9)          # Error / negative bg
MS_YELLOW = RGBColor(0xFF, 0xB9, 0x00)             # Accent yellow
MS_TEAL = RGBColor(0x03, 0x83, 0x87)               # Teal accent

# Typography Scale (pt) -- consistent hierarchy for all text
TEXT_DISPLAY = 46    # Hero / title slides
TEXT_H1 = 32         # Primary headings
TEXT_H2 = 28         # Slide titles (standard slides)
TEXT_H3 = 20         # Sub-headings / card headers
TEXT_BODY = 14       # Body copy
TEXT_BODY_SM = 12    # Secondary body / descriptions
TEXT_CAPTION = 10    # Captions, annotations, footnotes
TEXT_MICRO = 8       # Tiny labels, watermarks

# Spacing Grid (multiples of 0.08" base unit -- 8pt grid)
SPACE_XS = Inches(0.08)     # Tight: internal padding
SPACE_SM = Inches(0.15)     # Small: between related elements
SPACE_MD = Inches(0.25)     # Medium: standard spacing
SPACE_LG = Inches(0.4)      # Large: between sections
SPACE_XL = Inches(0.6)      # Extra large: major divisions
SPACE_2XL = Inches(0.8)     # 2X large: slide margins

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Content area (inside margins)
CONTENT_LEFT = Inches(0.8)
CONTENT_TOP = Inches(1.2)  # below title bar
CONTENT_WIDTH = Inches(11.0)
CONTENT_BOTTOM = Inches(6.8)  # above bottom bar

FONT_FAMILY = "Segoe UI"
FONT_MONO = "Cascadia Code"

# Asset paths
_ASSETS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
LOGO_PATH = os.path.join(_ASSETS_DIR, "microsoft.png")
LEAD_BG_PATH = os.path.join(_ASSETS_DIR, "lead-bg.jpg")

# ═════════════════════════════════════════════════════════════
# INITIALIZATION
# ═════════════════════════════════════════════════════════════

def create_presentation():
    """Create a new 16:9 Presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def new_blank_slide(prs):
    """Add a blank slide to the presentation."""
    return prs.slides.add_slide(prs.slide_layouts[6])


# ═════════════════════════════════════════════════════════════
# LOW-LEVEL SHAPE HELPERS
# ═════════════════════════════════════════════════════════════

def set_slide_bg(slide, color):
    """Set solid background color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_ms_logo(slide, left=None, top=None, width=Inches(1.4)):
    """Add the Microsoft logo watermark. Default: bottom-right with safe margins."""
    if not os.path.exists(LOGO_PATH):
        return None
    if left is None:
        left = SLIDE_WIDTH - width - Inches(0.5)
    if top is None:
        top = SLIDE_HEIGHT - Inches(0.65)
    try:
        return slide.shapes.add_picture(LOGO_PATH, left, top, width)
    except Exception:
        return None


def add_rect(slide, left, top, width, height, color, line_color=None):
    """Add a simple rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_top_accent_bar(slide, color=MS_BLUE, height=Inches(0.07)):
    """Add a colored accent bar at the top of a slide."""
    return add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, height, color)


def add_bottom_bar(slide, page_num=None, total=None):
    """Add page number at the bottom-left. No visible line - keeps slides clean."""
    if page_num is not None:
        tb = slide.shapes.add_textbox(CONTENT_LEFT, Inches(7.1), Inches(2), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{page_num}" + (f" / {total}" if total else "")
        run.font.size = Pt(10)
        run.font.color.rgb = MS_TEXT_MUTED
        run.font.name = FONT_FAMILY


def add_speaker_notes(slide, text):
    """Set speaker notes on a slide."""
    slide.notes_slide.notes_text_frame.text = text


def add_gradient_fill(shape, color_start, color_end, angle_deg=90):
    """Apply a linear gradient fill to any shape.

    Creates a smooth color transition from color_start to color_end.
    Use on rectangles, cards, or full-slide backgrounds for depth and atmosphere.

    angle_deg (PowerPoint convention):
        0   = left to right            90  = top to bottom (default)
        180 = right to left            270 = bottom to top
        45  = top-left to bottom-right 135 = top-right to bottom-left
        315 = bottom-left to top-right

    Example::
        bg = add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, MS_BLUE)
        add_gradient_fill(bg, MS_BLUE, MS_DARK_BLUE, angle_deg=135)
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    fill = shape.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = color_start
    stops[0].position = 0.0
    stops[1].color.rgb = color_end
    stops[1].position = 1.0

    # Set direction via XML (python-pptx doesn't expose gradient angle directly)
    gradFill = shape._element.spPr.find(qn('a:gradFill'))
    if gradFill is not None:
        for lin in gradFill.findall(qn('a:lin')):
            gradFill.remove(lin)
        lin = etree.SubElement(gradFill, qn('a:lin'))
        lin.set('ang', str(int(angle_deg * 60000)))
        lin.set('scaled', '1')


def add_gradient_fill_3(shape, color_start, color_mid, color_end, angle_deg=90):
    """Apply a three-stop linear gradient for richer color transitions.

    Example::
        add_gradient_fill_3(card, MS_BLUE, MS_BLUE_DARKER, MS_DARK_BLUE)
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    fill = shape.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = color_start
    stops[0].position = 0.0
    stops[1].color.rgb = color_mid
    stops[1].position = 0.5
    # Add a third stop
    from pptx.oxml import parse_xml
    gsLst = shape._element.spPr.find(qn('a:gradFill')).find(qn('a:gsLst'))
    gs3 = etree.SubElement(gsLst, qn('a:gs'))
    gs3.set('pos', '100000')
    srgb = etree.SubElement(gs3, qn('a:srgbClr'))
    srgb.set('val', '%02X%02X%02X' % (color_end[0], color_end[1], color_end[2]))

    gradFill = shape._element.spPr.find(qn('a:gradFill'))
    if gradFill is not None:
        for lin in gradFill.findall(qn('a:lin')):
            gradFill.remove(lin)
        lin = etree.SubElement(gradFill, qn('a:lin'))
        lin.set('ang', str(int(angle_deg * 60000)))
        lin.set('scaled', '1')


def add_shadow(shape, blur_pt=6, offset_pt=3, opacity=0.18, color="000000"):
    """Add a professional drop shadow to any shape.

    Adds depth and elevation, making cards and panels appear to float above
    the slide surface. The default values produce a subtle, sophisticated
    shadow matching Microsoft Fluent Design elevation levels.

    Presets (pass these values):
        Paper:   blur_pt=2,  offset_pt=1,  opacity=0.08  -- barely visible
        Subtle:  blur_pt=4,  offset_pt=2,  opacity=0.12  -- cards at rest
        Medium:  blur_pt=6,  offset_pt=3,  opacity=0.18  -- default, elevated
        Strong:  blur_pt=12, offset_pt=5,  opacity=0.25  -- modal / hero
        Deep:    blur_pt=20, offset_pt=8,  opacity=0.30  -- floating panel

    Example::
        card = add_rounded_card(slide, x, y, w, h)
        add_shadow(card)  # medium shadow
        add_shadow(card, blur_pt=4, offset_pt=2, opacity=0.12)  # subtle
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    spPr = shape._element.spPr

    # Remove any existing effects
    for ef in spPr.findall(qn('a:effectLst')):
        spPr.remove(ef)

    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))

    outerShdw.set('blurRad', str(int(blur_pt * 12700)))    # EMU conversion
    outerShdw.set('dist', str(int(offset_pt * 12700)))
    outerShdw.set('dir', '5400000')   # straight down (90 deg)
    outerShdw.set('rotWithShape', '0')

    srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgbClr.set('val', color)
    alpha = etree.SubElement(srgbClr, qn('a:alpha'))
    alpha.set('val', str(int(opacity * 100000)))
    return shape


def _lighten_color(color, factor=0.35):
    """Lighten an RGBColor by blending toward white.

    factor: 0.0 = unchanged, 1.0 = white. Default 0.35 = subtle lift.
    Returns a new RGBColor.
    """
    r = int(color[0] + (255 - color[0]) * factor)
    g = int(color[1] + (255 - color[1]) * factor)
    b = int(color[2] + (255 - color[2]) * factor)
    return RGBColor(r, g, b)


def _darken_color(color, factor=0.25):
    """Darken an RGBColor by blending toward black.

    factor: 0.0 = unchanged, 1.0 = black. Default 0.25 = subtle darken.
    Returns a new RGBColor.
    """
    r = int(color[0] * (1 - factor))
    g = int(color[1] * (1 - factor))
    b = int(color[2] * (1 - factor))
    return RGBColor(r, g, b)


# ═════════════════════════════════════════════════════════════
# RETURN TYPE & CONTRAST HELPERS
# ═════════════════════════════════════════════════════════════

# Consistent return type for composite elements.
# shape: the primary shape object
# left/top/width/height: bounding box in EMU for downstream stacking.
ElementBox = namedtuple("ElementBox", ["shape", "left", "top", "width", "height"])


def _luminance(color):
    """Calculate relative luminance (WCAG 2.0) of an RGBColor.

    Returns 0.0 (black) to 1.0 (white).
    """
    def _lin(c):
        s = c / 255.0
        return s / 12.92 if s <= 0.04045 else ((s + 0.055) / 1.055) ** 2.4
    return 0.2126 * _lin(color[0]) + 0.7152 * _lin(color[1]) + 0.0722 * _lin(color[2])


def auto_text_color(bg_color):
    """Return MS_WHITE or MS_DARK_BLUE based on background luminance.

    Ensures readable text regardless of background color.
    Replaces hardcoded light_fills tuple checks throughout the module.

    Example::
        text_c = auto_text_color(card_fill)
        add_textbox(slide, "Hello", ..., color=text_c)
    """
    return MS_DARK_BLUE if _luminance(bg_color) > 0.4 else MS_WHITE


def ensure_contrast(text_color, bg_color, min_ratio=3.0):
    """Return a corrected text_color if contrast ratio vs bg_color is too low.

    Uses WCAG 2.0 contrast ratio formula. Falls back to white-on-dark or
    dark-on-light when the original pairing is unreadable.

    Example::
        safe = ensure_contrast(MS_BLUE, MS_BLUE)  # -> MS_WHITE
        safe = ensure_contrast(MS_WHITE, MS_WHITE) # -> MS_DARK_BLUE
    """
    l1 = _luminance(text_color)
    l2 = _luminance(bg_color)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    ratio = (lighter + 0.05) / (darker + 0.05)
    if ratio >= min_ratio:
        return text_color
    return auto_text_color(bg_color)


# ═════════════════════════════════════════════════════════════
# SHAPE TEXT HELPERS
# ═════════════════════════════════════════════════════════════

def _set_shape_text(shape, text, font_size=14, color=MS_WHITE, bold=False,
                    alignment=PP_ALIGN.CENTER, v_align='middle',
                    font_name=FONT_FAMILY, italic=False, word_wrap=True,
                    shrink_to_fit=True,
                    margin_left=Inches(0.1), margin_right=Inches(0.1),
                    margin_top=Inches(0.05), margin_bottom=Inches(0.05)):
    """Set formatted text on a shape's native text_frame.

    Makes the text part of the shape itself -- movable as a single object in
    PowerPoint. Use instead of creating a separate ``add_textbox()`` overlay.

    v_align: 'top' | 'middle' | 'bottom' -- controls vertical anchoring.
    shrink_to_fit: when True, PowerPoint auto-shrinks text to prevent overflow.
    """
    from pptx.oxml.ns import qn

    tf = shape.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = int(margin_left)
    tf.margin_right = int(margin_right)
    tf.margin_top = int(margin_top)
    tf.margin_bottom = int(margin_bottom)

    if shrink_to_fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    anchor_map = {'top': 't', 'middle': 'ctr', 'bottom': 'b'}
    anchor_val = anchor_map.get(v_align, 'ctr')
    txBody = shape._element.txBody
    if txBody is not None:
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', anchor_val)

    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return shape


def _add_shape_paragraph(shape, text, font_size=14, color=MS_TEXT, bold=False,
                         alignment=PP_ALIGN.CENTER, font_name=FONT_FAMILY,
                         italic=False, space_before=None, space_after=None):
    """Add an additional paragraph to a shape that already has text.

    Use after ``_set_shape_text()`` to add a second line (e.g., label below
    a metric value). Returns the shape for chaining.
    """
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.alignment = alignment
    if space_before is not None:
        p.space_before = space_before
    if space_after is not None:
        p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return shape


def group_shapes(slide, shapes_list):
    """Group multiple shapes into a single selectable/movable unit.

    Returns the GroupShape. All shapes in shapes_list must belong to slide.
    Grouped shapes can be ungrouped in PowerPoint for editing.

    Example::
        card = add_rounded_card(slide, x, y, w, h)
        icon = add_icon_circle(slide, ix, iy, size, color, "1")
        group = group_shapes(slide, [card, icon])
    """
    from pptx.oxml.ns import qn
    from lxml import etree
    from copy import deepcopy

    spTree = slide.shapes._spTree

    # Calculate bounding box of all shapes
    min_l = min(int(s.left) for s in shapes_list)
    min_t = min(int(s.top) for s in shapes_list)
    max_r = max(int(s.left) + int(s.width) for s in shapes_list)
    max_b = max(int(s.top) + int(s.height) for s in shapes_list)

    # Build group shape XML
    grpSp = etree.SubElement(spTree, qn('p:grpSp'))
    nvGrpSpPr = etree.SubElement(grpSp, qn('p:nvGrpSpPr'))
    cNvPr = etree.SubElement(nvGrpSpPr, qn('p:cNvPr'))
    cNvPr.set('id', str(len(spTree) + 100))
    cNvPr.set('name', 'Group')
    etree.SubElement(nvGrpSpPr, qn('p:cNvGrpSpPr'))
    etree.SubElement(nvGrpSpPr, qn('p:nvPr'))

    grpSpPr = etree.SubElement(grpSp, qn('p:grpSpPr'))
    xfrm = etree.SubElement(grpSpPr, qn('a:xfrm'))
    off = etree.SubElement(xfrm, qn('a:off'))
    off.set('x', str(min_l))
    off.set('y', str(min_t))
    ext = etree.SubElement(xfrm, qn('a:ext'))
    ext.set('cx', str(max_r - min_l))
    ext.set('cy', str(max_b - min_t))
    chOff = etree.SubElement(xfrm, qn('a:chOff'))
    chOff.set('x', str(min_l))
    chOff.set('y', str(min_t))
    chExt = etree.SubElement(xfrm, qn('a:chExt'))
    chExt.set('cx', str(max_r - min_l))
    chExt.set('cy', str(max_b - min_t))

    # Move each shape's XML element into the group
    for shape in shapes_list:
        el = shape._element
        spTree.remove(el)
        grpSp.append(el)

    return grpSp


# ═════════════════════════════════════════════════════════════
# TEXT HELPERS
# ═════════════════════════════════════════════════════════════

def add_textbox(slide, text, left, top, width, height,
                font_size=16, color=MS_TEXT, bold=False,
                alignment=PP_ALIGN.LEFT, font_name=FONT_FAMILY,
                italic=False, shrink_to_fit=False, v_align=None):
    """Add a styled textbox. Returns the textbox shape.

    shrink_to_fit: when True, PowerPoint auto-shrinks text to prevent overflow.
    v_align: 'top' | 'middle' | 'bottom' (default: None = PowerPoint default).
    """
    from pptx.oxml.ns import qn

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if shrink_to_fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if v_align:
        anchor_map = {'top': 't', 'middle': 'ctr', 'bottom': 'b'}
        bodyPr = tf._txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', anchor_map.get(v_align, 't'))
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return tb


def add_rich_text(slide, parts, left, top, width, height,
                  alignment=PP_ALIGN.LEFT):
    """Add text with mixed formatting. parts = [(text, {font_size, color, bold, ...}), ...]"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    for text, fmt in parts:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fmt.get("font_size", 16))
        run.font.color.rgb = fmt.get("color", MS_TEXT)
        run.font.bold = fmt.get("bold", False)
        run.font.italic = fmt.get("italic", False)
        run.font.name = fmt.get("font_name", FONT_FAMILY)
    return tb


def add_bullet_list(slide, items, left, top, width, height=Inches(5),
                    font_size=14, color=MS_TEXT, spacing=Pt(8)):
    """Add a bullet list. items = [str] or [(bold_prefix, rest_text)]."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        if isinstance(item, tuple):
            bold_part, rest_part = item
            r1 = p.add_run()
            r1.text = bold_part
            r1.font.size = Pt(font_size)
            r1.font.color.rgb = MS_DARK_BLUE
            r1.font.bold = True
            r1.font.name = FONT_FAMILY
            r2 = p.add_run()
            r2.text = rest_part
            r2.font.size = Pt(font_size)
            r2.font.color.rgb = color
            r2.font.name = FONT_FAMILY
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(font_size)
            r.font.color.rgb = color
            r.font.name = FONT_FAMILY
    return tb


def add_gradient_textbox(slide, text, left, top, width, height,
                         color_start, color_end, angle_deg=0,
                         font_size=16, bold=False,
                         alignment=PP_ALIGN.LEFT, font_name=FONT_FAMILY,
                         italic=False):
    """Add a textbox where the text itself has a gradient fill.

    This is a native PowerPoint feature where characters are filled with a
    smooth color transition rather than a solid color. Creates a polished,
    premium look for titles and headings.

    color_start, color_end: RGBColor pair defining the gradient.
    angle_deg: 0=left-to-right (default), 90=top-to-bottom, etc.

    Example::
        add_gradient_textbox(slide, "Key Insights", x, y, w, h,
                             MS_BLUE, MS_LIGHT_BLUE, font_size=28, bold=True)
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name

    # Apply gradient fill on text characters.
    # OOXML schema requires fill elements (gradFill) BEFORE font elements
    # (latin, ea, cs). python-pptx appends latin at the end, so we must
    # build the rPr children in the correct order.
    rPr = run._r.get_or_add_rPr()

    # Remove any existing fill
    for sf in rPr.findall(qn('a:solidFill')):
        rPr.remove(sf)

    # Build gradFill element
    gradFill = etree.Element(qn('a:gradFill'))
    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

    gs1 = etree.SubElement(gsLst, qn('a:gs'))
    gs1.set('pos', '0')
    srgb1 = etree.SubElement(gs1, qn('a:srgbClr'))
    srgb1.set('val', '%02X%02X%02X' % (color_start[0], color_start[1], color_start[2]))

    gs2 = etree.SubElement(gsLst, qn('a:gs'))
    gs2.set('pos', '100000')
    srgb2 = etree.SubElement(gs2, qn('a:srgbClr'))
    srgb2.set('val', '%02X%02X%02X' % (color_end[0], color_end[1], color_end[2]))

    lin = etree.SubElement(gradFill, qn('a:lin'))
    lin.set('ang', str(int(angle_deg * 60000)))
    lin.set('scaled', '1')

    # Insert gradFill BEFORE the first font element (a:latin, a:ea, a:cs, a:sym)
    font_tags = [qn('a:latin'), qn('a:ea'), qn('a:cs'), qn('a:sym'),
                 qn('a:hlinkClick'), qn('a:hlinkMouseOver'), qn('a:rtl'), qn('a:extLst')]
    inserted = False
    for child in list(rPr):
        if child.tag in font_tags:
            rPr.insert(list(rPr).index(child), gradFill)
            inserted = True
            break
    if not inserted:
        rPr.append(gradFill)

    return tb


# ═════════════════════════════════════════════════════════════
# SHAPE HELPERS
# ═════════════════════════════════════════════════════════════

def add_icon_circle(slide, left, top, size, color, text="", text_color=None):
    """Add a colored circle with optional centered text.

    Args:
        text_color: RGBColor for the text. Defaults to MS_WHITE when the
                    circle fill is dark, or MS_DARK_BLUE when the fill is
                    a light/white color. Pass explicitly to override.
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        # Auto-detect text color using luminance when not provided
        if text_color is None:
            text_color = auto_text_color(color)
        _set_shape_text(shape, text,
                        font_size=int(size / Pt(1) * 0.4),
                        color=text_color, bold=True,
                        alignment=PP_ALIGN.CENTER, v_align='middle',
                        word_wrap=False,
                        margin_left=0, margin_right=0,
                        margin_top=0, margin_bottom=0)
    return shape


def add_rounded_card(slide, left, top, width, height, fill=MS_WHITE, border=None,
                     corner_radius=0.05):
    """Add a rounded rectangle card with subtle corner rounding.

    Args:
        corner_radius: 0.0 = sharp corners, 0.5 = pill. Default 0.05 = subtle rounding.
                       Use 0.0 for dense/tabular layouts. Never use default OOXML rounding
                       (~0.13) -- it looks cartoonish at small sizes.
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.0)
    else:
        shape.line.fill.background()
    # Subtle rounding (python-pptx adjustment: 0.0 sharp to 0.5 full-pill)
    try:
        shape.adjustments[0] = corner_radius
    except Exception:
        pass
    return shape


def add_arrow_right(slide, left, top, width=Inches(0.5), height=Inches(0.4), color=MS_BLUE):
    """Add a right-pointing chevron arrow. Use between horizontal flow cards."""
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_arrow_down(slide, left, top, width=Inches(0.4), height=Inches(0.5), color=MS_BLUE):
    """Add a downward-pointing chevron arrow for vertical flow diagrams.

    The CHEVRON shape points right by default; rotating 90 degrees makes it
    point down. After rotation the visual footprint is approximately height
    wide x width tall, so for a visually square arrow keep width == height.

    Positioning tip: to center the arrow between two vertically-stacked cards::

        add_arrow_down(slide,
            left=card_left + (card_width - Inches(0.4)) / 2,
            top=bottom_of_upper_card + Inches(0.1))
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.rotation = 90  # 90 deg clockwise = points down
    return shape


def add_arrow_up(slide, left, top, width=Inches(0.4), height=Inches(0.5), color=MS_BLUE):
    """Add an upward-pointing chevron arrow for vertical flow diagrams.

    See add_arrow_down() for positioning notes -- same logic applies.
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.rotation = 270  # 270 deg clockwise = points up
    return shape


def add_header_card(slide, left, top, width, height, header_text, color,
                    header_height=Inches(0.5)):
    """Add a rounded card with a colored header banner and white body.

    Layers a colored rounded rect (full height) behind a white rounded rect
    (offset below the header), creating a seamless colored header with rounded
    top corners. Returns (header_shape, body_shape) for further content placement.

    Content area starts at: top + header_height + Inches(0.1)
    """
    # 1. Full-height colored card (peek shows as header + provides outer border)
    header_shape = add_rounded_card(slide, left, top, width, height, fill=color, border=color)
    # 2. White body card overlaid, offset below header (no border -- the outer
    #    card's border already frames the card; a body border creates an ugly
    #    straight line at the header/body junction)
    body_top = top + header_height
    body_height = height - header_height
    body_shape = add_rounded_card(slide, left, body_top, width, body_height,
                                  fill=MS_WHITE, border=None, corner_radius=0.0)
    # 3. Header text
    add_textbox(slide, header_text, left + Inches(0.15), top + Inches(0.05),
                width - Inches(0.3), header_height,
                font_size=15, color=MS_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    return header_shape, body_shape


def add_elevated_card(slide, left, top, width, height, fill=MS_WHITE, border=None,
                      corner_radius=0.05, shadow="medium"):
    """Add a rounded card with a drop shadow for visual depth.

    Combines add_rounded_card with add_shadow for easy elevation.
    Use instead of add_rounded_card when the card needs to visually "pop".

    shadow: "paper" | "subtle" | "medium" | "strong" | "deep" | None
    """
    card = add_rounded_card(slide, left, top, width, height, fill, border, corner_radius)
    if shadow:
        presets = {
            "paper":  (2, 1, 0.08),
            "subtle": (4, 2, 0.12),
            "medium": (6, 3, 0.18),
            "strong": (12, 5, 0.25),
            "deep":   (20, 8, 0.30),
        }
        blur, offset, opacity = presets.get(shadow, presets["medium"])
        add_shadow(card, blur_pt=blur, offset_pt=offset, opacity=opacity)
    return card


def add_badge(slide, text, left, top, bg_color=MS_BLUE, text_color=MS_WHITE,
              font_size=9, width=None, height=Inches(0.28)):
    """Add a pill-shaped status badge (e.g., "GA", "Preview", "New").

    Perfect for tagging features, products, or status on cards and slides.
    Width auto-sizes to text when None.

    Common patterns::
        add_badge(slide, "GA",         x, y, bg_color=MS_GREEN)
        add_badge(slide, "Preview",    x, y, bg_color=MS_ORANGE)
        add_badge(slide, "New",        x, y, bg_color=MS_BLUE)
        add_badge(slide, "Deprecated", x, y, bg_color=MS_RED)
    """
    if width is None:
        char_width = font_size * 0.6 / 72.0
        width = Inches(len(text) * char_width + 0.25)

    text_color = ensure_contrast(text_color, bg_color)
    pill = add_rounded_card(slide, left, top, width, height,
                            fill=bg_color, border=None, corner_radius=0.5)
    _set_shape_text(pill, text, font_size=font_size, color=text_color, bold=True,
                    alignment=PP_ALIGN.CENTER, v_align='middle',
                    margin_left=Inches(0.05), margin_right=Inches(0.05),
                    margin_top=0, margin_bottom=0)
    return pill


def add_divider_line(slide, left, top, width, color=MS_MID_GRAY, thickness=Pt(1)):
    """Add a subtle horizontal divider line.

    Use between content sections on the same slide for visual separation
    without creating a new section. Much lighter than a full-width bar.
    """
    return add_rect(slide, left, top, width, thickness, color)


def add_progress_bar(slide, left, top, width, height=Inches(0.18),
                     progress=0.7, bar_color=MS_BLUE, track_color=MS_LIGHT_GRAY,
                     label="", show_pct=True):
    """Add a horizontal progress/completion bar.

    progress: 0.0 to 1.0 (percentage complete).
    Optionally shows percentage text and/or label.

    Example::
        add_progress_bar(slide, x, y, Inches(4), progress=0.85, label="Adoption")
    """
    # Track (background)
    track = add_rounded_card(slide, left, top, width, height,
                             fill=track_color, border=None, corner_radius=0.5)
    # Fill (foreground)
    fill_w = max(int(width * progress), Inches(0.08))
    fill_bar = add_rounded_card(slide, left, top, fill_w, height,
                                fill=bar_color, border=None, corner_radius=0.5)
    # Label (above bar)
    if label:
        add_textbox(slide, label, left, top - Inches(0.22), width, Inches(0.2),
                    font_size=TEXT_CAPTION, color=MS_TEXT_MUTED, bold=True)
    # Percentage (right of bar)
    if show_pct:
        pct_text = f"{int(progress * 100)}%"
        add_textbox(slide, pct_text, left + width + Inches(0.1), top - Inches(0.01),
                    Inches(0.5), height,
                    font_size=TEXT_CAPTION, color=MS_TEXT, bold=True)
    return fill_bar


def add_checklist(slide, items, left, top, width, height=Inches(5),
                  font_size=13, check_color=MS_GREEN, spacing=Pt(10)):
    """Add a checklist with checkmark symbols.

    items = [str] for all checked, or [(str, bool)] for mixed checked/unchecked.

    The checkmark style uses Unicode symbols with color coding:
        Checked:   green checkmark
        Unchecked: gray circle
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing

        # Determine check state
        if isinstance(item, tuple):
            text, checked = item
        else:
            text, checked = item, True

        # Checkmark/circle symbol
        r_icon = p.add_run()
        r_icon.text = "\u2713 " if checked else "\u25CB "
        r_icon.font.size = Pt(font_size)
        r_icon.font.color.rgb = check_color if checked else MS_MID_GRAY
        r_icon.font.bold = True
        r_icon.font.name = FONT_FAMILY

        # Item text
        r_text = p.add_run()
        r_text.text = text
        r_text.font.size = Pt(font_size)
        r_text.font.color.rgb = MS_TEXT if checked else MS_TEXT_MUTED
        r_text.font.name = FONT_FAMILY
    return tb


def add_gradient_card(slide, left, top, width, height, color_start, color_end,
                      angle_deg=90, corner_radius=0.05, shadow="subtle"):
    """Add a rounded card with a gradient fill and optional shadow.

    Combines add_rounded_card + add_gradient_fill + add_shadow for a single call.
    Great for hero cards, section headers, and call-to-action panels.

    Example::
        add_gradient_card(slide, x, y, w, h, MS_BLUE, MS_DARK_BLUE, angle_deg=135)
    """
    card = add_rounded_card(slide, left, top, width, height,
                            fill=color_start, border=None, corner_radius=corner_radius)
    add_gradient_fill(card, color_start, color_end, angle_deg)
    if shadow:
        presets = {
            "paper":  (2, 1, 0.08),
            "subtle": (4, 2, 0.12),
            "medium": (6, 3, 0.18),
            "strong": (12, 5, 0.25),
            "deep":   (20, 8, 0.30),
        }
        blur, offset, opacity = presets.get(shadow, presets["medium"])
        add_shadow(card, blur_pt=blur, offset_pt=offset, opacity=opacity)
    return card


# ═════════════════════════════════════════════════════════════
# LAYOUT & LAYER UTILITIES
# ═════════════════════════════════════════════════════════════

# LAYER ORDER (z-order = insertion order; last shape added renders on top):
#   1. Background fills / full-slide images / color bands
#   2. Structural containers (cards, panels) -- draw ALL containers first
#   3. Arrows / connectors                  -- draw AFTER all cards
#   4. Text boxes / icons / number badges   -- draw LAST; always on top
#
# ALWAYS use two passes in loops:
#   Pass 1: for item in items: draw_card(...)
#   Pass 2: for arrow in arrows: draw_arrow(...)
#   Pass 3: for item in items: draw_text(...)


def estimate_text_height(text, font_size_pt, width_inches, padding_inches=0.22,
                          line_spacing_factor=1.45, min_lines=1):
    """Estimate the Inches height needed to display text at font_size_pt in width_inches.

    Uses a character-width heuristic (Segoe UI) with a safety multiplier.
    Returns an Inches() EMU value ready to pass directly to shape constructors.

    Changes from earlier version:
        - Wider character estimate (0.58 vs 0.50) to avoid underflow.
        - min_lines parameter to guarantee minimum height.
        - Handles long unbreakable tokens (URLs, identifiers).
        - 10% safety margin on final result.

    Example::
        h = estimate_text_height(long_note, 13, 10.0)
        add_callout_box(slide, long_note, left, top, Inches(10), h)
    """
    # Average character width in inches for Segoe UI (conservative)
    avg_char_w = (font_size_pt * 0.58) / 72.0
    chars_per_line = max(10, int(width_inches / avg_char_w))

    # Word-wrap simulation (handles long tokens)
    words = text.split()
    lines, line_len = 1, 0
    for word in words:
        needed = len(word) + (1 if line_len else 0)
        if line_len + needed > chars_per_line:
            # Long unbreakable word: count wrapped lines it would need
            if len(word) > chars_per_line:
                extra_lines = (len(word) - 1) // chars_per_line
                lines += extra_lines + 1
                line_len = len(word) % chars_per_line or chars_per_line
            else:
                lines += 1
                line_len = len(word)
        else:
            line_len += needed

    lines = max(lines, min_lines)
    line_h = (font_size_pt * line_spacing_factor) / 72.0
    raw = lines * line_h + padding_inches
    return Inches(raw * 1.10)  # 10% safety margin


# ═════════════════════════════════════════════════════════════
# COMPOSITE CONTENT ELEMENTS
# ═════════════════════════════════════════════════════════════

def add_callout_box(slide, text, left, top, width, height=None,
                    bg=MS_CALLOUT_BG, accent=MS_BLUE, font_size=13):
    """Add a callout box with a left accent bar, auto-sized to content.

    When height is None (default) the box height is calculated from the text
    length, font_size, and width so the box is exactly as tall as it needs to
    be -- no awkward empty space. Pass an explicit Inches() to force a fixed
    height.

    Returns the actual height used so callers can chain vertical positions::

        h = add_callout_box(slide, text1, left, top, width)
        add_callout_box(slide, text2, left, top + h + Inches(0.15), width)
    """
    if height is None:
        height = estimate_text_height(
            text, font_size,
            float(width) / 914400 - 0.45,  # account for 0.3+0.15 text inset
            padding_inches=0.25,
        )
    # Background with embedded text (single selectable unit in PowerPoint)
    card = add_rounded_card(slide, left, top, width, height, fill=bg, border=None,
                            corner_radius=0.04)
    add_rect(slide, left, top, Pt(5), height, accent)
    _set_shape_text(card, text, font_size=font_size, color=MS_DARK_BLUE, bold=True,
                    alignment=PP_ALIGN.LEFT, v_align='top',
                    margin_left=Inches(0.3), margin_right=Inches(0.15),
                    margin_top=Inches(0.1), margin_bottom=Inches(0.05))
    return height


def add_warning_box(slide, text, left, top, width, height=None):
    """Add an orange warning callout box. height=None auto-sizes to content."""
    return add_callout_box(slide, text, left, top, width, height,
                           bg=RGBColor(0xFF, 0xF4, 0xE5), accent=MS_ORANGE)


def add_code_block(slide, code, left, top, width, height):
    """Add a dark-themed code block with a monospace font and blue left border."""
    # Dark background with embedded code text (single unit in PowerPoint)
    bg = add_rect(slide, left, top, width, height, MS_CODE_BG)
    bg.shadow.inherit = False
    # Blue left accent
    add_rect(slide, left, top, Pt(4), height, MS_BLUE)
    # Code text embedded in background shape
    _set_shape_text(bg, code, font_size=10, color=MS_CODE_TEXT, bold=False,
                    font_name=FONT_MONO, alignment=PP_ALIGN.LEFT, v_align='top',
                    margin_left=Inches(0.2), margin_right=Inches(0.1),
                    margin_top=Inches(0.1), margin_bottom=Inches(0.05))
    return bg


def add_styled_table(slide, data, left, top, width, col_widths=None,
                     header_color=MS_DARK_BLUE, font_size=12):
    """Add a professionally styled table with alternating rows."""
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.38 * rows))
    table = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            if isinstance(w, float) and w < 1.0:
                # Treat as a ratio of total width
                table.columns[i].width = int(width * w)
            elif isinstance(w, float):
                table.columns[i].width = int(Inches(w))
            else:
                table.columns[i].width = w
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.name = FONT_FAMILY
            run.font.size = Pt(font_size if ri > 0 else font_size + 1)
            p.alignment = PP_ALIGN.LEFT
            if ri == 0:
                run.font.bold = True
                run.font.color.rgb = MS_WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                run.font.color.rgb = MS_TEXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = MS_CALLOUT_BG if ri % 2 == 0 else MS_WHITE
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
    return ts


def add_metric_card(slide, metric, label, x, y, w=Inches(3.5), h=Inches(2.5),
                    color=MS_BLUE, sublabel="", trend="", trend_positive=True):
    """Add a single big-number metric card.

    Text is embedded in the card shape (single selectable object in PowerPoint).
    Internal positions scale proportionally with card height.

    trend: optional trend text, e.g. "+12% YoY". Shown below the label with
           a green up-arrow or red down-arrow based on trend_positive.

    Example::
        add_metric_card(slide, "98.5%", "Uptime SLA", x, y,
                        trend="+0.3%", trend_positive=True)
    """
    card = add_elevated_card(slide, x, y, w, h, fill=MS_WHITE,
                             border=color, shadow="subtle")
    add_rect(slide, x, y, w, Inches(0.06), color)
    # All text embedded in card with proportional spacing
    h_emu = int(h)
    _set_shape_text(card, str(metric),
                    font_size=40, color=color, bold=True,
                    alignment=PP_ALIGN.CENTER, v_align='middle',
                    margin_left=Inches(0.1), margin_right=Inches(0.1),
                    margin_top=int(h_emu * 0.06), margin_bottom=int(h_emu * 0.08))
    _add_shape_paragraph(card, label,
                         font_size=14, color=MS_DARK_BLUE,
                         alignment=PP_ALIGN.CENTER,
                         space_before=Pt(6))
    if sublabel:
        _add_shape_paragraph(card, sublabel,
                             font_size=11, color=MS_TEXT_MUTED,
                             alignment=PP_ALIGN.CENTER,
                             space_before=Pt(4))
    if trend:
        arrow = "\u25B2 " if trend_positive else "\u25BC "
        t_color = MS_GREEN if trend_positive else MS_RED
        _add_shape_paragraph(card, arrow + trend,
                             font_size=TEXT_BODY_SM, color=t_color, bold=True,
                             alignment=PP_ALIGN.CENTER,
                             space_before=Pt(4))
    return card


def add_numbered_items(slide, items, left, top, width, item_height=Inches(1.1),
                       colors=None):
    """Add numbered items with colored circles. items = [(title, desc)]."""
    if colors is None:
        colors = [MS_BLUE, MS_DARK_BLUE, MS_GREEN, MS_ORANGE, MS_PURPLE] * 3
    for i, (title, desc) in enumerate(items):
        y = top + i * item_height
        c = colors[i % len(colors)]
        bg_color = MS_LIGHT_GRAY if i % 2 == 0 else MS_WHITE
        ih = int(item_height)
        card_h = item_height - Inches(0.05)
        add_rounded_card(slide, left, y, width, card_h,
                         fill=bg_color, border=MS_MID_GRAY)
        add_rect(slide, left, y, Pt(5), card_h, c)
        # Proportional positions (scale with item_height)
        add_icon_circle(slide, left + Inches(0.2), y + int(ih * 0.08),
                        Inches(0.5), c, str(i + 1))
        add_textbox(slide, title, left + Inches(0.9), y + int(ih * 0.03),
                    width - Inches(1.1), int(ih * 0.35),
                    font_size=15, color=MS_DARK_BLUE, bold=True)
        add_textbox(slide, desc, left + Inches(0.9), y + int(ih * 0.38),
                    width - Inches(1.1), int(ih * 0.55),
                    font_size=12, color=MS_TEXT_MUTED)


def add_card_grid(slide, cards, left, top, cols=2, card_w=Inches(5.5),
                  card_h=Inches(2.3), gap_x=Inches(0.35), gap_y=Inches(0.35)):
    """Add cards in a grid. cards = [(color, title, desc), ...]."""
    for i, (color, title, desc) in enumerate(cards):
        col = i % cols
        row = i // cols
        x = left + col * (card_w + gap_x)
        y = top + row * (card_h + gap_y)
        ch = int(card_h)  # EMU for proportional offsets
        add_rounded_card(slide, x, y, card_w, card_h, fill=color, border=MS_MID_GRAY)
        add_rounded_card(slide, x, y + Inches(0.08), card_w, card_h - Inches(0.08),
                         fill=MS_WHITE, border=MS_MID_GRAY)
        add_icon_circle(slide, x + Inches(0.25), y + int(ch * 0.13), Inches(0.5), color, str(i + 1))
        add_textbox(slide, title, x + Inches(0.95), y + int(ch * 0.10),
                    card_w - Inches(1.2), int(ch * 0.18),
                    font_size=16, color=MS_DARK_BLUE, bold=True)
        add_textbox(slide, desc, x + Inches(0.95), y + int(ch * 0.30),
                    card_w - Inches(1.2), int(ch * 0.65),
                    font_size=12, color=MS_TEXT_MUTED)


def add_pillar_cards(slide, pillars, left=CONTENT_LEFT, top=CONTENT_TOP, height=Inches(5.2),
                     min_gap=Inches(0.15)):
    """Add vertical pillar cards with proportional internal layout.

    All internal offsets (circle, title, description) are derived from ``height``
    so the cards never overflow regardless of how tall or short they are.

    pillars = [(color, num, title, desc), ...].
    """
    n = len(pillars)
    card_w = (CONTENT_WIDTH - (n - 1) * min_gap) / n
    h = float(height)  # EMU
    # Proportional breakpoints (fraction of total card height)
    circle_r   = Inches(0.4)
    circle_top = top + h * 0.07
    title_top  = top + h * 0.30
    title_h    = h * 0.18
    desc_top   = top + h * 0.52
    desc_h     = h * 0.44
    for i, (color, num, title_text, desc) in enumerate(pillars):
        x = left + i * (card_w + min_gap)
        add_rounded_card(slide, x, top, card_w, height, fill=color, border=color)
        add_rounded_card(slide, x, top + Inches(0.08), card_w, height - Inches(0.08),
                         fill=MS_WHITE, border=color)
        add_icon_circle(slide, x + card_w / 2 - circle_r, circle_top, circle_r * 2, color, num)
        add_textbox(slide, title_text, x + Inches(0.1), title_top,
                    card_w - Inches(0.2), title_h,
                    font_size=15, color=MS_DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, desc, x + Inches(0.1), desc_top,
                    card_w - Inches(0.2), desc_h,
                    font_size=11, color=MS_TEXT_MUTED, alignment=PP_ALIGN.CENTER)


def add_quote_block(slide, quote, attribution="", left=CONTENT_LEFT, top=Inches(2.0),
                    width=Inches(10), accent_color=MS_BLUE):
    """Add a large quotation block with left accent bar and attribution.

    Professional style for customer testimonials, key statements, or
    expert opinions. Uses oversized quote marks and italic text.

    Example::
        add_quote_block(slide,
            "GitHub Copilot reduced our onboarding time by 55%.",
            "- VP Engineering, Contoso Ltd")
    """
    q_height = estimate_text_height(quote, 20, float(width) / 914400 - 1.0)
    total_h = q_height + (Inches(0.5) if attribution else Inches(0.1))

    # White card with shadow for elevation
    card = add_elevated_card(slide, left, top, width, total_h + Inches(0.4),
                             fill=MS_WHITE, border=MS_LIGHT_GRAY, shadow="subtle")
    # Left accent bar
    add_rect(slide, left, top, Pt(5), total_h + Inches(0.4), accent_color)
    # Opening quote mark (large, decorative)
    add_textbox(slide, "\u201C", left + Inches(0.3), top - Inches(0.05),
                Inches(0.6), Inches(0.7),
                font_size=48, color=accent_color, bold=True)
    # Quote text
    add_textbox(slide, quote, left + Inches(0.8), top + Inches(0.15),
                width - Inches(1.2), q_height,
                font_size=20, color=MS_DARK_BLUE, italic=True)
    # Attribution
    if attribution:
        add_textbox(slide, attribution, left + Inches(0.8),
                    top + q_height + Inches(0.15),
                    width - Inches(1.2), Inches(0.35),
                    font_size=TEXT_BODY_SM, color=MS_TEXT_MUTED, bold=True)
    return card


def add_stats_row(slide, stats, left=CONTENT_LEFT, top=Inches(1.8),
                  width=None, card_h=Inches(1.6), gap=Inches(0.25)):
    """Add a horizontal row of big-number statistics.

    stats = [(value, label), ...] or [(value, label, color), ...]

    Each stat is displayed as a large bold number with a small label below.
    Cards are evenly distributed across the available width.

    Example::
        add_stats_row(slide, [
            ("55%", "Faster Onboarding"),
            ("3.2x", "Code Review Speed"),
            ("40%", "Less Boilerplate"),
        ])
    """
    if width is None:
        width = CONTENT_WIDTH
    n = len(stats)
    card_w = (width - (n - 1) * gap) / n

    cards = []
    for i, stat in enumerate(stats):
        value = stat[0]
        label = stat[1]
        color = stat[2] if len(stat) > 2 else MS_BLUE

        x = left + i * (card_w + gap)

        # Card with top accent and shadow
        card = add_elevated_card(slide, x, top, card_w, card_h,
                                 fill=MS_WHITE, border=MS_LIGHT_GRAY, shadow="subtle")
        add_rect(slide, x, top, card_w, Inches(0.05), color)

        # Text embedded in card with proportional spacing
        ch = int(card_h)
        _set_shape_text(card, str(value),
                        font_size=36, color=color, bold=True,
                        alignment=PP_ALIGN.CENTER, v_align='middle',
                        margin_left=SPACE_SM, margin_right=SPACE_SM,
                        margin_top=int(ch * 0.06), margin_bottom=int(ch * 0.08))
        _add_shape_paragraph(card, label,
                             font_size=TEXT_BODY_SM, color=MS_DARK_BLUE,
                             alignment=PP_ALIGN.CENTER,
                             space_before=Pt(4))
        cards.append(card)
    return cards


def add_kpi_card(slide, value, label, trend="", trend_positive=True,
                 x=CONTENT_LEFT, y=Inches(2.0), w=Inches(2.8), h=Inches(2.2),
                 color=MS_BLUE):
    """Alias for add_metric_card with trend support. Prefer add_metric_card directly."""
    return add_metric_card(slide, value, label, x, y, w, h, color,
                           trend=trend, trend_positive=trend_positive)


def add_comparison_columns(slide, col_left, col_right,
                           left=CONTENT_LEFT, top=CONTENT_TOP,
                           width=None, height=Inches(5.0),
                           left_title="Before", right_title="After",
                           left_color=MS_MID_GRAY, right_color=MS_BLUE):
    """Add a side-by-side comparison layout (Before/After, Option A/B, etc.).

    col_left  = [str, ...] -- bullet items for the left column.
    col_right = [str, ...] -- bullet items for the right column.

    Creates two cards with contrasting headers and bullet content.

    Example::
        add_comparison_columns(slide,
            ["Manual deployments", "Slow reviews", "No standards"],
            ["CI/CD pipelines", "AI-assisted reviews", "Copilot coding standards"],
            left_title="Without Copilot", right_title="With Copilot")
    """
    if width is None:
        width = CONTENT_WIDTH
    gap = Inches(0.3)
    col_w = (width - gap) / 2

    # Left column
    lx = left
    add_header_card_with_bullets(slide, left_title, col_left,
                                lx, top, col_w, height,
                                header_color=left_color, font_size=TEXT_BODY_SM)

    # Right column
    rx = left + col_w + gap
    add_header_card_with_bullets(slide, right_title, col_right,
                                rx, top, col_w, height,
                                header_color=right_color, font_size=TEXT_BODY_SM)

    # VS divider circle
    vs_x = left + col_w + gap / 2 - Inches(0.25)
    vs_y = top + Inches(0.1)
    add_icon_circle(slide, vs_x, vs_y, Inches(0.5), MS_DARK_BLUE, "VS")


def add_feature_grid(slide, features, left=CONTENT_LEFT, top=CONTENT_TOP,
                     cols=3, card_w=None, card_h=Inches(1.8),
                     gap=Inches(0.2)):
    """Add a grid of feature cards with numbered circles and descriptions.

    features = [(title, description), ...] or [(title, description, color), ...]

    Each card has a numbered circle, bold title, and description text.
    Cards are arranged in a responsive grid.

    Example::
        add_feature_grid(slide, [
            ("Code Completion", "AI-powered suggestions as you type"),
            ("Chat", "Natural language coding assistance"),
            ("CLI", "Terminal-based AI help"),
        ])
    """
    if card_w is None:
        card_w = (CONTENT_WIDTH - (cols - 1) * gap) / cols

    default_colors = [MS_BLUE, MS_DARK_BLUE, MS_GREEN, MS_ORANGE, MS_PURPLE, MS_RED]

    for i, feat in enumerate(features):
        title = feat[0]
        desc = feat[1]
        color = feat[2] if len(feat) > 2 else default_colors[i % len(default_colors)]

        col = i % cols
        row = i // cols
        x = left + col * (card_w + gap)
        y = top + row * (card_h + gap)
        ch = int(card_h)  # EMU for proportional offsets

        # Card with subtle shadow
        card = add_elevated_card(slide, x, y, card_w, card_h,
                                 fill=MS_WHITE, border=MS_LIGHT_GRAY, shadow="paper")
        # Top accent
        add_rect(slide, x, y, card_w, Inches(0.04), color)
        # Number circle
        add_icon_circle(slide, x + Inches(0.15), y + int(ch * 0.10),
                        Inches(0.4), color, str(i + 1))
        # Title
        add_textbox(slide, title, x + Inches(0.65), y + int(ch * 0.08),
                    card_w - Inches(0.85), int(ch * 0.22),
                    font_size=TEXT_BODY, color=MS_DARK_BLUE, bold=True)
        # Description
        add_textbox(slide, desc, x + Inches(0.15), y + int(ch * 0.38),
                    card_w - Inches(0.3), int(ch * 0.55),
                    font_size=TEXT_BODY_SM, color=MS_TEXT_MUTED)


def add_colored_columns(slide, columns, left=CONTENT_LEFT, top=CONTENT_TOP,
                        width=None, gap=Inches(0.3), title_font_size=TEXT_H3,
                        body_font_size=TEXT_BODY_SM, bullet_symbol="\u00B7 ",
                        gradient_titles=True):
    """Add side-by-side columns with colored bold titles and bullet lists.

    A lightweight alternative to add_header_card_with_bullets -- no card
    background, no header banner, just colored title text with plain bullets.

    columns = [(title, items_list, color), ...]  -- solid color title
         or   [(title, items_list, (color_start, color_end)), ...]  -- gradient title

    gradient_titles: if True (default) and color is a single RGBColor, a lighter
        variant is automatically derived for a subtle gradient. Pass a tuple of
        two RGBColor values for explicit control.

    Example::
        add_colored_columns(slide, [
            ("Scaling with the customer", [
                "Consciously planning adoption",
                "Setting R&R for Microsoft AND the customer",
            ], MS_BLUE),
            ("Strong IP Developed", [
                "Playbook with detailed guidance",
                "Hands-on workshops for multiple languages/IDEs",
            ], MS_GREEN),
        ])
    """
    if width is None:
        width = CONTENT_WIDTH
    n = len(columns)
    col_w = (width - (n - 1) * gap) / n

    for i, (title, items, color) in enumerate(columns):
        x = left + i * (col_w + gap)

        # Determine gradient colors
        if isinstance(color, tuple) and len(color) == 2:
            c_start, c_end = color
        elif gradient_titles:
            c_start = color
            c_end = _lighten_color(color, 0.35)
        else:
            c_start, c_end = color, None

        # Title with gradient or solid text
        if c_end is not None:
            add_gradient_textbox(slide, title, x, top, col_w, Inches(0.6),
                                 c_start, c_end,
                                 font_size=title_font_size, bold=True)
        else:
            add_textbox(slide, title, x, top, col_w, Inches(0.6),
                        font_size=title_font_size, color=c_start, bold=True)

        # Bullet list
        bullets_top = top + Inches(0.65)
        bullets_h = Inches(4.5)
        tb = slide.shapes.add_textbox(x, bullets_top, col_w, bullets_h)
        tf = tb.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"{bullet_symbol}{item}"
            run.font.size = Pt(body_font_size)
            run.font.color.rgb = MS_TEXT
            run.font.name = FONT_FAMILY


def add_layered_architecture(slide, layers, left=CONTENT_LEFT, top=Inches(1.5),
                             width=None, layer_h=Inches(0.9), gap=Inches(0.08)):
    """Add a horizontal layered architecture diagram (stacked bars).

    layers = [(label, color), ...] ordered top-to-bottom.
    Creates a classic architecture stack diagram where each layer spans
    the full width, perfect for showing technology stacks, OSI layers,
    or platform architectures.

    Example::
        add_layered_architecture(slide, [
            ("Application Layer", MS_BLUE),
            ("API Gateway / Middleware", MS_DARK_BLUE),
            ("Platform Services", MS_GREEN),
            ("Infrastructure (AKS / VMs)", MS_DARK_GRAY),
        ])
    """
    if width is None:
        width = CONTENT_WIDTH

    for i, (label, color) in enumerate(layers):
        y = top + i * (layer_h + gap)
        # Layer bar with subtle shadow
        bar = add_elevated_card(slide, left, y, width, layer_h,
                                fill=color, border=None,
                                corner_radius=0.03, shadow="paper")
        # Label text embedded in bar (auto contrast)
        text_color = auto_text_color(color)
        _set_shape_text(bar, label, font_size=TEXT_H3, color=text_color, bold=True,
                        alignment=PP_ALIGN.CENTER, v_align='middle',
                        margin_left=Inches(0.3), margin_right=Inches(0.3),
                        margin_top=Inches(0.05), margin_bottom=Inches(0.05))


def add_agenda_list(slide, items, left=CONTENT_LEFT, top=Inches(1.5),
                    width=None, highlight_index=None):
    """Add a numbered agenda/table-of-contents list.

    items = ["Introduction", "Architecture Overview", "Demo", ...]
    highlight_index: index of the currently active item (0-based).

    Active item is highlighted with blue background; others are subtle.
    Perfect for agenda slides and section markers showing progress.

    Example::
        add_agenda_list(slide, [
            "Introduction & Goals",
            "Architecture Deep Dive",
            "Live Demo",
            "Q&A",
        ], highlight_index=1)
    """
    if width is None:
        width = CONTENT_WIDTH
    item_h = Inches(0.7)
    gap = Inches(0.08)

    for i, item in enumerate(items):
        y = top + i * (item_h + gap)
        is_active = (i == highlight_index)

        # Background card
        bg_color = MS_BLUE if is_active else MS_LIGHT_GRAY
        text_color = MS_WHITE if is_active else MS_TEXT
        num_bg = MS_DARK_BLUE if is_active else MS_MID_GRAY

        card = add_rounded_card(slide, left, y, width, item_h,
                                fill=bg_color, border=None, corner_radius=0.04)
        if is_active:
            add_shadow(card, blur_pt=4, offset_pt=2, opacity=0.15)

        # Number circle
        add_icon_circle(slide, left + Inches(0.15), y + Inches(0.1),
                        Inches(0.5), num_bg, str(i + 1))
        # Text
        add_textbox(slide, item, left + Inches(0.8), y + Inches(0.1),
                    width - Inches(1.0), item_h - Inches(0.2),
                    font_size=TEXT_H3, color=text_color, bold=is_active)


def add_icon_row(slide, items, left=CONTENT_LEFT, top=Inches(2.5),
                 width=None, icon_size=Inches(0.6)):
    """Add a horizontal row of icon circles with labels below.

    items = [(symbol, label), ...] or [(symbol, label, color), ...]

    Creates a row of colored circles with centered icons and text labels,
    evenly spaced across the width. Good for capability overviews.

    Example::
        add_icon_row(slide, [
            ("\u2699", "Settings"),
            ("\u2605", "Favorites"),
            ("\u2709", "Messages"),
        ])
    """
    if width is None:
        width = CONTENT_WIDTH
    n = len(items)
    slot_w = width / n
    default_colors = [MS_BLUE, MS_DARK_BLUE, MS_GREEN, MS_ORANGE, MS_PURPLE, MS_RED]

    for i, item in enumerate(items):
        symbol = item[0]
        label = item[1]
        color = item[2] if len(item) > 2 else default_colors[i % len(default_colors)]

        cx = left + i * slot_w + (slot_w - icon_size) / 2
        add_icon_circle(slide, cx, top, icon_size, color, symbol)
        add_textbox(slide, label, left + i * slot_w, top + icon_size + Inches(0.1),
                    slot_w, Inches(0.5),
                    font_size=TEXT_BODY_SM, color=MS_DARK_BLUE, bold=True,
                    alignment=PP_ALIGN.CENTER)


def add_pricing_table(slide, tiers, left=CONTENT_LEFT, top=Inches(1.3),
                      width=None, height=Inches(5.5), highlight_index=None):
    """Add a pricing/comparison table with tiered columns.

    tiers = [(name, price, features_list, color), ...]
    highlight_index: index of the recommended tier (gets elevated + badge).

    Each tier is a card with header, price, and feature checklist.

    Example::
        add_pricing_table(slide, [
            ("Basic", "$10/mo", ["5 users", "10GB storage"], MS_MID_GRAY),
            ("Pro", "$25/mo", ["25 users", "100GB", "Priority support"], MS_BLUE),
            ("Enterprise", "Custom", ["Unlimited", "1TB", "Dedicated CSM"], MS_DARK_BLUE),
        ], highlight_index=1)
    """
    if width is None:
        width = CONTENT_WIDTH
    n = len(tiers)
    gap = Inches(0.2)
    col_w = (width - (n - 1) * gap) / n

    for i, (name, price, features, color) in enumerate(tiers):
        x = left + i * (col_w + gap)
        is_featured = (i == highlight_index)

        # Card (featured gets elevation)
        y_offset = Inches(-0.15) if is_featured else 0
        card_h = height + (Inches(0.15) if is_featured else 0)

        card = add_elevated_card(slide, x, top + y_offset, col_w, card_h,
                                 fill=MS_WHITE, border=color,
                                 shadow="strong" if is_featured else "subtle")

        # Header
        add_rounded_card(slide, x, top + y_offset, col_w, Inches(0.55),
                         fill=color, border=None, corner_radius=0.0)
        add_textbox(slide, name, x + SPACE_SM, top + y_offset + Inches(0.08),
                    col_w - SPACE_MD, Inches(0.4),
                    font_size=TEXT_H3, color=MS_WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # "Recommended" badge
        if is_featured:
            badge_w = Inches(1.2)
            add_badge(slide, "Recommended", x + (col_w - badge_w) / 2,
                      top + y_offset + Inches(0.63),
                      bg_color=MS_ORANGE, width=badge_w)

        # Price
        price_top = top + y_offset + Inches(0.75) + (Inches(0.35) if is_featured else 0)
        add_textbox(slide, price, x + SPACE_SM, price_top,
                    col_w - SPACE_MD, Inches(0.6),
                    font_size=28, color=color, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # Features checklist
        feat_top = price_top + Inches(0.65)
        add_checklist(slide, features, x + Inches(0.2), feat_top,
                      col_w - Inches(0.4), font_size=TEXT_BODY_SM,
                      check_color=color)


def add_swot_grid(slide, strengths, weaknesses, opportunities, threats,
                  left=CONTENT_LEFT, top=CONTENT_TOP, width=None, height=Inches(5.0)):
    """Add a SWOT analysis grid (2x2 matrix).

    Each quadrant is a colored card with a header and bullet list.
    Standard consulting layout for strategic analysis.

    Example::
        add_swot_grid(slide,
            ["Strong brand", "Market leader"],
            ["High costs", "Legacy systems"],
            ["Cloud migration", "AI adoption"],
            ["New competitors", "Regulation"])
    """
    if width is None:
        width = CONTENT_WIDTH
    gap = Inches(0.15)
    cell_w = (width - gap) / 2
    cell_h = (height - gap) / 2

    quadrants = [
        ("Strengths",     strengths,     MS_GREEN,     left,              top),
        ("Weaknesses",    weaknesses,    MS_ORANGE,    left + cell_w + gap, top),
        ("Opportunities", opportunities, MS_BLUE,      left,              top + cell_h + gap),
        ("Threats",        threats,       MS_RED,       left + cell_w + gap, top + cell_h + gap),
    ]
    for title, items, color, qx, qy in quadrants:
        add_header_card_with_bullets(slide, title, items,
                                    qx, qy, cell_w, cell_h,
                                    header_color=color, font_size=TEXT_BODY_SM)


def add_maturity_model(slide, levels, current_level=None,
                       left=CONTENT_LEFT, top=Inches(2.0),
                       width=None, height=Inches(4.0)):
    """Add a maturity model visualization (staircase pattern).

    levels = [(label, description), ...] ordered from lowest to highest.
    current_level: 0-based index of the current maturity level (highlighted).

    Creates an ascending staircase of cards, with the current level highlighted.
    Standard consulting pattern for capability assessments.

    Example::
        add_maturity_model(slide, [
            ("Initial", "Ad-hoc processes"),
            ("Managed", "Basic planning"),
            ("Defined", "Standardized processes"),
            ("Measured", "Data-driven"),
            ("Optimized", "Continuous improvement"),
        ], current_level=2)
    """
    if width is None:
        width = CONTENT_WIDTH
    n = len(levels)
    step_w = width / n
    step_base_h = height * 0.25  # minimum height for lowest step
    step_increment = (height - step_base_h) / max(n - 1, 1)

    for i, (label, desc) in enumerate(levels):
        x = left + i * step_w
        step_h = step_base_h + i * step_increment
        y = top + height - step_h
        is_current = (i == current_level)

        color = MS_BLUE if is_current else MS_LIGHT_GRAY
        text_c = MS_WHITE if is_current else MS_DARK_BLUE
        desc_c = RGBColor(0xDD, 0xEE, 0xFF) if is_current else MS_TEXT_MUTED

        card = add_elevated_card(slide, x + Inches(0.04), y, step_w - Inches(0.08),
                                 step_h, fill=color, border=MS_MID_GRAY,
                                 shadow="medium" if is_current else "paper")
        # Level number
        add_textbox(slide, f"L{i + 1}", x + Inches(0.1), y + Inches(0.08),
                    step_w - Inches(0.16), Inches(0.3),
                    font_size=TEXT_CAPTION, color=desc_c, bold=True,
                    alignment=PP_ALIGN.CENTER)
        # Label
        add_textbox(slide, label, x + Inches(0.1), y + Inches(0.35),
                    step_w - Inches(0.16), Inches(0.4),
                    font_size=TEXT_BODY, color=text_c, bold=True,
                    alignment=PP_ALIGN.CENTER)
        # Description (only if step is tall enough)
        if step_h > Inches(1.2):
            add_textbox(slide, desc, x + Inches(0.1), y + Inches(0.8),
                        step_w - Inches(0.16), step_h - Inches(1.0),
                        font_size=TEXT_CAPTION, color=desc_c,
                        alignment=PP_ALIGN.CENTER)


def add_roadmap(slide, phases, left=CONTENT_LEFT, top=Inches(2.0),
                width=None, phase_h=Inches(1.2)):
    """Add a horizontal phased roadmap with milestone markers.

    phases = [(phase_label, items_list, color), ...]
    Each phase is a card with a header and bullet items inside.

    More detailed than add_timeline -- each phase has room for multiple
    items, making it suitable for project roadmaps and release planning.

    Example::
        add_roadmap(slide, [
            ("Phase 1 - Foundation", ["Set up infra", "Core APIs"], MS_BLUE),
            ("Phase 2 - Build", ["UI development", "Integrations"], MS_GREEN),
            ("Phase 3 - Launch", ["Beta testing", "GA release"], MS_ORANGE),
        ])
    """
    if width is None:
        width = CONTENT_WIDTH
    n = len(phases)
    gap = Inches(0.15)
    phase_w = (width - (n - 1) * gap) / n

    for i, (label, items, color) in enumerate(phases):
        x = left + i * (phase_w + gap)
        items_h = Inches(0.3) * len(items) + Inches(0.2)
        total_h = max(phase_h, Inches(0.6) + items_h)

        add_header_card_with_bullets(slide, label, items,
                                    x, top, phase_w, total_h,
                                    header_color=color, font_size=TEXT_BODY_SM)

        # Arrow connector
        if i < n - 1:
            arrow_x = x + phase_w + Inches(0.02)
            arrow_y = top + total_h / 2 - Inches(0.2)
            add_arrow_right(slide, arrow_x, arrow_y,
                            width=gap - Inches(0.04), height=Inches(0.4),
                            color=color)


# ═════════════════════════════════════════════════════════════
# SLIDE TEMPLATES
# ═════════════════════════════════════════════════════════════

def create_standard_slide(prs, title, page_num=None, total=None, notes=""):
    """Create a standard content slide with title, blue underline, logo, bottom bar.

    No top accent bar -- the title underline is the sole brand accent. This gives
    content slides a cleaner, more editorial feel and avoids the heavy banner that
    competes visually with the main content area.

    Returns the slide for adding content below the title (CONTENT_TOP = 1.2").
    """
    slide = new_blank_slide(prs)
    set_slide_bg(slide, MS_WHITE)
    # Title (moved up slightly since there is no top bar)
    add_textbox(slide, title, CONTENT_LEFT, Inches(0.2), CONTENT_WIDTH, Inches(0.7),
                font_size=28, color=MS_DARK_BLUE, bold=True)
    # Blue underline -- the primary brand accent on content slides
    add_rect(slide, CONTENT_LEFT, Inches(0.95), CONTENT_WIDTH, Pt(3), MS_BLUE)
    add_bottom_bar(slide, page_num, total)
    add_ms_logo(slide)
    if notes:
        add_speaker_notes(slide, notes)
    return slide


def create_section_divider(prs, title, subtitle="", notes=""):
    """Create a section divider with MS_DARK_BLUE background, decorative shapes, and logo.
    Uses dark navy (#243A5E) so the Microsoft logo is clearly readable."""
    slide = new_blank_slide(prs)
    set_slide_bg(slide, MS_DARK_BLUE)
    # Left accent bar (bright blue on dark background)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, MS_BLUE)
    # Decorative circles (use slightly lighter shade for contrast on dark bg)
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1.2), Inches(4.5), Inches(4.5))
    c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x2E, 0x4A, 0x6E); c1.line.fill.background()
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(5.8), Inches(1.5), Inches(1.5))
    c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x2E, 0x4A, 0x6E); c2.line.fill.background()
    # Title
    add_textbox(slide, title, Inches(1.0), Inches(2.5), Inches(10), Inches(1.2),
                font_size=44, color=MS_WHITE, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, Inches(1.0), Inches(4.1), Inches(9), Inches(0.8),
                    font_size=18, color=MS_ACCENT_LIGHT)
    add_ms_logo(slide, left=SLIDE_WIDTH - Inches(1.9), top=SLIDE_HEIGHT - Inches(0.7), width=Inches(1.5))
    if notes:
        add_speaker_notes(slide, notes)
    return slide


def create_lead_slide(prs, title, subtitle="", meta="", notes="",
                      level="", use_bg_image=True):
    """Create a title/lead slide with background image, blue band, and logo.

    Args:
        level: Content level shown in the circle badge.
               Must be "L100", "L200", "L300", or "L400".
               When provided, the circle badge displays the level
               with a "Content Level" sublabel.
    """
    slide = new_blank_slide(prs)
    set_slide_bg(slide, MS_WHITE)
    # Background image
    if use_bg_image and os.path.exists(LEAD_BG_PATH):
        slide.shapes.add_picture(LEAD_BG_PATH, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    # Blue band on right
    add_rect(slide, Inches(8.5), Inches(0), Inches(4.833), SLIDE_HEIGHT, MS_BLUE)

    # Level badge in decorative circles
    if level:
        c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(1.5), Inches(4.5), Inches(4.5))
        c1.fill.solid(); c1.fill.fore_color.rgb = MS_DARK_BLUE; c1.line.fill.background()
        c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(2.2), Inches(3.1), Inches(3.1))
        c2.fill.solid(); c2.fill.fore_color.rgb = MS_BLUE; c2.line.fill.background()

        # Large level text
        tf = c2.text_frame; tf.word_wrap = True
        tf.paragraphs[0].space_after = Pt(0)
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        run1 = p1.add_run()
        run1.text = level.upper()
        run1.font.size = Pt(44)
        run1.font.color.rgb = MS_WHITE
        run1.font.bold = True
        run1.font.name = FONT_FAMILY
        # Sublabel
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(2)
        run2 = p2.add_run()
        run2.text = "Content Level"
        run2.font.size = Pt(11)
        run2.font.color.rgb = MS_ACCENT_LIGHT
        run2.font.bold = False
        run2.font.name = FONT_FAMILY

    # Logo top-left
    add_ms_logo(slide, left=Inches(0.6), top=Inches(0.5), width=Inches(1.6))
    # Title
    add_textbox(slide, title, Inches(0.8), Inches(2.0), Inches(7), Inches(1.8),
                font_size=46, color=MS_DARK_BLUE, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.8), Inches(4.2), Inches(7), Inches(0.6),
                    font_size=20, color=MS_BLUE)
    if meta:
        add_textbox(slide, meta, Inches(0.8), Inches(5.2), Inches(7), Inches(0.4),
                    font_size=13, color=MS_TEXT_MUTED)
    if notes:
        add_speaker_notes(slide, notes)
    return slide


def create_closing_slide(prs, title="Key Takeaways", takeaways=None,
                         cta_title="", cta_url="", cta_items=None,
                         page_num=None, total=None, notes=""):
    """Create a split closing slide: blue left (takeaways) + white right (CTA)."""
    slide = new_blank_slide(prs)
    set_slide_bg(slide, MS_WHITE)
    # Blue band left
    add_rect(slide, Inches(0), Inches(0), Inches(5.2), SLIDE_HEIGHT, MS_BLUE)
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(4.8), Inches(3.5), Inches(3.5))
    c.fill.solid(); c.fill.fore_color.rgb = MS_DARK_BLUE; c.line.fill.background()
    # Left title
    add_textbox(slide, title, Inches(0.6), Inches(0.6), Inches(4), Inches(0.6),
                font_size=30, color=MS_WHITE, bold=True)
    # Blue underline below title
    add_rect(slide, Inches(0.6), Inches(1.2), Inches(2.2), Pt(3), MS_WHITE)
    # Takeaways
    if takeaways:
        for i, t in enumerate(takeaways):
            y = Inches(1.5) + i * Inches(1.3)
            add_icon_circle(slide, Inches(0.6), y, Inches(0.4), MS_DARK_BLUE, str(i + 1))
            add_textbox(slide, t, Inches(1.2), y - Inches(0.05), Inches(3.5), Inches(1.1),
                        font_size=12, color=RGBColor(0xDD, 0xEE, 0xFF))
    add_ms_logo(slide, left=Inches(1.5), top=Inches(6.5), width=Inches(1.3))
    # Right side
    if cta_title:
        add_textbox(slide, cta_title, Inches(5.8), Inches(1.2), Inches(6.5), Inches(1.2),
                    font_size=32, color=MS_DARK_BLUE, bold=True)
    if cta_url:
        add_textbox(slide, cta_url, Inches(5.8), Inches(2.5), Inches(6.5), Inches(0.35),
                    font_size=14, color=MS_BLUE)
    if cta_items:
        for i, item in enumerate(cta_items):
            if isinstance(item, str):
                big, small = item, ""
            else:
                big, small = item[0], item[1] if len(item) > 1 else ""
            x = Inches(5.8) + Inches(i * 2.35)
            add_rounded_card(slide, x, Inches(3.1), Inches(2.1), Inches(1.7),
                             fill=MS_WHITE, border=MS_BLUE)
            add_textbox(slide, big, x + Inches(0.1), Inches(3.2), Inches(1.9), Inches(0.8),
                        font_size=18, color=MS_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
            add_textbox(slide, small, x + Inches(0.1), Inches(4.05), Inches(1.9), Inches(0.4),
                        font_size=11, color=MS_TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, "Questions?", Inches(5.8), Inches(5.5), Inches(6.5), Inches(0.5),
                font_size=26, color=MS_DARK_BLUE, bold=True)
    add_bottom_bar(slide, page_num, total)
    if notes:
        add_speaker_notes(slide, notes)
    return slide


def create_two_tone_slide(prs, title, left_color=MS_DARK_BLUE, right_color=MS_WHITE,
                          split_inches=5.0, page_num=None, total=None, notes=""):
    """Create a split-background slide: colored left panel + white right content area.

    The left panel occupies split_inches of the slide and uses a dark
    background with white title text. The right panel is white for content.

    Perfect for:
        - Key message + supporting detail
        - Quote/stat + explanation
        - Navigation sidebar + content

    Returns (slide, left_width, right_left, right_width) for content placement.
    """
    slide = new_blank_slide(prs)
    set_slide_bg(slide, right_color)

    left_w = Inches(split_inches)

    # Left panel
    add_rect(slide, Inches(0), Inches(0), left_w, SLIDE_HEIGHT, left_color)

    # Decorative circle (subtle depth on dark panel)
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(5.0),
                               Inches(3.5), Inches(3.5))
    c.fill.solid()
    c.fill.fore_color.rgb = MS_NAVY_LIGHT
    c.line.fill.background()

    # Title on left panel
    add_textbox(slide, title, Inches(0.5), Inches(0.5),
                left_w - Inches(0.8), Inches(1.0),
                font_size=TEXT_H1, color=MS_WHITE, bold=True)
    # White underline
    add_rect(slide, Inches(0.5), Inches(1.55), Inches(2.0), Pt(3), MS_WHITE)

    add_ms_logo(slide, left=Inches(0.5), top=SLIDE_HEIGHT - Inches(0.7), width=Inches(1.3))
    add_bottom_bar(slide, page_num, total)
    if notes:
        add_speaker_notes(slide, notes)

    right_left = Inches(split_inches + 0.4)
    right_width = Inches(13.333 - split_inches - 0.8)

    return slide, left_w, right_left, right_width


def create_gradient_slide(prs, title, color_start=MS_BLUE, color_end=MS_DARK_BLUE,
                          angle_deg=135, page_num=None, total=None, notes=""):
    """Create a slide with a full-bleed gradient background.

    All text appears in white on the gradient. Use for impact statements,
    section intros, or key data slides that need visual emphasis.

    Returns the slide for adding content (use white text colors).
    """
    slide = new_blank_slide(prs)

    # Full-slide gradient background rectangle
    bg = add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, color_start)
    add_gradient_fill(bg, color_start, color_end, angle_deg)

    # Title
    add_textbox(slide, title, CONTENT_LEFT, Inches(0.4), CONTENT_WIDTH, Inches(0.8),
                font_size=TEXT_H1, color=MS_WHITE, bold=True)
    # White underline
    add_rect(slide, CONTENT_LEFT, Inches(1.2), Inches(3.0), Pt(3), MS_WHITE)

    add_ms_logo(slide)
    add_bottom_bar(slide, page_num, total)
    if notes:
        add_speaker_notes(slide, notes)
    return slide


def create_agenda_slide(prs, title, items, highlight_index=None,
                        page_num=None, total=None, notes=""):
    """Create a full agenda slide with numbered items and optional highlighting.

    Convenience wrapper that creates a standard slide and populates it
    with add_agenda_list.

    Example::
        create_agenda_slide(prs, "Agenda", [
            "Welcome & Introductions",
            "Architecture Overview",
            "Live Demo",
            "Next Steps & Q&A",
        ], highlight_index=1)
    """
    slide = create_standard_slide(prs, title, page_num, total, notes)
    add_agenda_list(slide, items, highlight_index=highlight_index)
    return slide


def create_impact_slide(prs, headline, subtext="", stats=None,
                        page_num=None, total=None, notes=""):
    """Create a high-impact statement slide with gradient background.

    Features a large centered headline with optional subtitle and
    bottom stats row. Use for key takeaways or dramatic data points.

    Example::
        create_impact_slide(prs, "55% Faster Onboarding",
            "Across 180+ development teams at Contoso",
            stats=[("55%", "Faster"), ("3.2x", "More PRs"), ("40%", "Less Bugs")])
    """
    slide = new_blank_slide(prs)
    bg = add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, MS_BLUE)
    add_gradient_fill(bg, MS_DARK_BLUE, MS_BLUE, angle_deg=135)

    # Headline (centered, large)
    add_textbox(slide, headline, Inches(1.0), Inches(2.0),
                Inches(11.333), Inches(2.0),
                font_size=TEXT_DISPLAY, color=MS_WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    if subtext:
        add_textbox(slide, subtext, Inches(1.5), Inches(4.2),
                    Inches(10.333), Inches(0.8),
                    font_size=TEXT_H3, color=MS_ACCENT_LIGHT,
                    alignment=PP_ALIGN.CENTER)

    # Stats row at bottom (on semi-transparent cards)
    if stats:
        n = len(stats)
        card_w = Inches(2.5)
        total_w = n * card_w + (n - 1) * Inches(0.3)
        start_x = (SLIDE_WIDTH - total_w) / 2
        for i, stat in enumerate(stats):
            val = stat[0]
            lbl = stat[1]
            sx = start_x + i * (card_w + Inches(0.3))
            sy = Inches(5.2)
            # Dark card on gradient
            add_rounded_card(slide, sx, sy, card_w, Inches(1.4),
                             fill=MS_DARK_BLUE, border=None, corner_radius=0.05)
            add_textbox(slide, str(val), sx + Inches(0.1), sy + Inches(0.1),
                        card_w - Inches(0.2), Inches(0.7),
                        font_size=30, color=MS_WHITE, bold=True,
                        alignment=PP_ALIGN.CENTER)
            add_textbox(slide, lbl, sx + Inches(0.1), sy + Inches(0.8),
                        card_w - Inches(0.2), Inches(0.4),
                        font_size=TEXT_BODY_SM, color=MS_ACCENT_LIGHT,
                        alignment=PP_ALIGN.CENTER)

    add_ms_logo(slide)
    add_bottom_bar(slide, page_num, total)
    if notes:
        add_speaker_notes(slide, notes)
    return slide


# ═════════════════════════════════════════════════════════════
# MICROSOFT-STYLE COMPONENTS
# ═════════════════════════════════════════════════════════════

def add_title_icon_badge(slide, left, top, symbol="\u2693", size=Inches(0.55),
                         bg_color=MS_BLUE, text_color=MS_WHITE):
    """Add a circular icon badge next to a slide title (e.g., person, gear symbol).

    Common Unicode symbols that render well in Segoe UI:
        \\u2699  Gear           \\u2696  Scales
        \\u2605  Star           \\u2630  Trigram (menu)
        \\u2602  Umbrella       \\u2622  Radioactive
        \\u25B6  Play           \\u2714  Checkmark
        \\u2764  Heart          \\u26A1  Lightning
        \\u2B50  Star           \\u2709  Envelope

    Tip: place at (title_right + Inches(0.15), title_top + Inches(0.05)).
    Returns the circle shape.
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = symbol
    run.font.size = Pt(int(float(size) / 914400 * 72 * 0.45))
    run.font.color.rgb = text_color
    run.font.bold = True
    run.font.name = FONT_FAMILY
    return shape


def add_blue_speech_panel(slide, text, left, top, width, height,
                          bg_color=MS_BLUE, text_color=MS_WHITE,
                          font_size=13, accent_bar=True):
    """Add a solid blue panel with white text and left accent bar.

    Inspired by the blue "thesis statement" boxes in Microsoft adoption slides
    (screenshots 1 & 3). The panel uses a filled rounded rectangle with optional
    darker left accent strip.

    Unlike add_callout_box (light background, dark text), this is a
    high-contrast branded panel for key statements.

    Returns the background shape.
    """
    text_color = ensure_contrast(text_color, bg_color)
    card = add_rounded_card(slide, left, top, width, height,
                            fill=bg_color, border=None, corner_radius=0.04)
    if accent_bar:
        add_rect(slide, left, top, Pt(5), height, MS_DARK_BLUE)
    _set_shape_text(card, text, font_size=font_size, color=text_color, bold=False,
                    alignment=PP_ALIGN.LEFT, v_align='middle',
                    margin_left=Inches(0.25), margin_right=Inches(0.15),
                    margin_top=Inches(0.15), margin_bottom=Inches(0.1))
    return card


def add_header_card_with_bullets(slide, header_text, bullets, left, top, width, height,
                                 header_color=MS_BLUE, header_height=Inches(0.5),
                                 font_size=11, bullet_symbol="\u2022 "):
    """Add a card with a colored header banner and bulleted body content.

    Extends the add_header_card pattern with built-in bullet list rendering.
    Matches the style from "Measure Adoption" slide (screenshot 3).

    bullets = [str, ...] -- each string is one bullet point.
    Returns (header_shape, body_shape, textbox).

    Example::
        add_header_card_with_bullets(
            slide, "Organization",
            ["Acceptance rate", "Total active users"],
            x, y, Inches(2.2), Inches(3.0))
    """
    header_shape, body_shape = add_header_card(
        slide, left, top, width, height, header_text, header_color,
        header_height=header_height)

    # Build bullet text inside body
    body_top = top + header_height + Inches(0.12)
    body_h = height - header_height - Inches(0.2)
    tb = slide.shapes.add_textbox(left + Inches(0.15), body_top,
                                  width - Inches(0.3), body_h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"{bullet_symbol}{item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = MS_TEXT
        run.font.name = FONT_FAMILY
    return header_shape, body_shape, tb


def add_timeline(slide, phases, left=CONTENT_LEFT, top=Inches(1.3),
                 box_height=Inches(0.55), desc_height=Inches(1.8),
                 available_width=None, bar_color=MS_BLUE):
    """Add a horizontal timeline/roadmap with milestone boxes and descriptions.

    Matches the adoption plan timeline style (screenshot 2).

    phases = [(label, week_label, description), ...]
        label:       text inside the blue milestone box (e.g., "Kick Off")
        week_label:  text below the connector line (e.g., "Week 1 & 2")
        description: body text below the week label

    Layout (top-down):
        [ Blue Box ] --- dotted line --- [ Blue Box ]
          Week 1 & 2                      Week 3
          Description                     Description

    Returns list of milestone box shapes for further customization.
    """
    if available_width is None:
        available_width = CONTENT_WIDTH
    n = len(phases)
    gap = Inches(0.15)
    box_w = (available_width - (n - 1) * gap) / n
    connector_y = top + box_height + Inches(0.08)
    week_y = connector_y + Inches(0.05)
    desc_y = week_y + Inches(0.35)
    shapes = []
    for i, (label, week_label, description) in enumerate(phases):
        x = left + i * (box_w + gap)

        # Milestone box (text embedded in shape)
        box = add_rounded_card(slide, x, top, box_w, box_height,
                               fill=bar_color, border=None, corner_radius=0.06)
        _set_shape_text(box, label, font_size=11, color=MS_WHITE, bold=True,
                        alignment=PP_ALIGN.CENTER, v_align='middle',
                        margin_left=Inches(0.05), margin_right=Inches(0.05),
                        margin_top=Inches(0.03), margin_bottom=Inches(0.03))
        shapes.append(box)

        # Dotted connector line between boxes
        if i < n - 1:
            line_left = x + box_w
            line_right = line_left + gap
            mid_y = connector_y
            conn = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          line_left, mid_y, gap, Pt(2))
            conn.fill.solid()
            conn.fill.fore_color.rgb = MS_MID_GRAY
            conn.line.fill.background()

        # Week label
        add_textbox(slide, week_label, x, week_y, box_w, Inches(0.3),
                    font_size=10, color=MS_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, description, x, desc_y, box_w, desc_height,
                    font_size=10, color=MS_TEXT_MUTED)

    # Arrow at end suggesting continuation
    arrow_x = left + n * (box_w + gap) - gap + Inches(0.05)
    if arrow_x + Inches(0.35) < left + available_width:
        add_textbox(slide, "\u2192", arrow_x, top + Inches(0.05),
                    Inches(0.4), box_height,
                    font_size=20, color=MS_BLUE, bold=True)

    return shapes


def add_activity_bars(slide, activities, left=CONTENT_LEFT, top=Inches(5.8),
                      width=None, bar_height=Inches(0.32), gap=Inches(0.05),
                      colors=None):
    """Add horizontal activity/workstream bars at the bottom of a slide.

    Matches the continuous bars in the timeline slide (screenshot 2) that span
    the full width and represent ongoing activities (surveys, comms, etc.).

    activities = ["Surveys", "Plan adjustment ...", "Communications ..."]
    colors: list of RGBColor or None for auto-gradient from dark to light blue.

    Returns list of bar shapes.
    """
    if width is None:
        width = CONTENT_WIDTH
    if colors is None:
        colors = [MS_DARK_BLUE, MS_BLUE, MS_BLUE_DARKER, MS_GREEN, MS_PURPLE]
    bars = []
    for i, text in enumerate(activities):
        y = top + i * (bar_height + gap)
        c = colors[i % len(colors)]
        bar = add_rounded_card(slide, left, y, width, bar_height,
                               fill=c, border=None, corner_radius=0.08)
        _set_shape_text(bar, text, font_size=10, color=auto_text_color(c), bold=True,
                        alignment=PP_ALIGN.CENTER, v_align='middle',
                        margin_left=Inches(0.2), margin_right=Inches(0.2),
                        margin_top=Inches(0.02), margin_bottom=Inches(0.02))
        bars.append(bar)
    return bars


def add_process_flow(slide, steps, left=CONTENT_LEFT, top=Inches(2.0),
                     box_w=Inches(2.5), box_h=Inches(1.4),
                     arrow_w=Inches(0.6), colors=None,
                     annotations=None):
    """Add a horizontal process flow diagram with boxes and block arrows.

    Matches the adoption strategy flow (screenshot 4):
    [Hackathon] --> [Train the Trainers] --> [Scale-up]

    steps = ["Hackathon", "Train the Trainers", "Scale-up"]
    colors: list of RGBColor per step, or None for automatic dark-to-light.
    annotations: optional list of (step_index, text) for bullets below a step.

    Returns list of box shapes.
    """
    n = len(steps)
    if colors is None:
        palette = [MS_DARK_BLUE, MS_BLUE, MS_LIGHT_BLUE]
        colors = [palette[i % len(palette)] for i in range(n)]

    total_w = n * box_w + (n - 1) * arrow_w
    start_x = left + (CONTENT_WIDTH - total_w) / 2  # center the flow

    boxes = []
    for i, label in enumerate(steps):
        x = start_x + i * (box_w + arrow_w)
        text_color = MS_WHITE if colors[i] != MS_LIGHT_BLUE else MS_DARK_BLUE

        # Box (text embedded in shape)
        box = add_rounded_card(slide, x, top, box_w, box_h,
                               fill=colors[i], border=None, corner_radius=0.06)
        _set_shape_text(box, label, font_size=18, color=text_color, bold=True,
                        alignment=PP_ALIGN.CENTER, v_align='middle',
                        margin_left=Inches(0.1), margin_right=Inches(0.1),
                        margin_top=Inches(0.1), margin_bottom=Inches(0.1))
        boxes.append(box)

        # Arrow between steps
        if i < n - 1:
            arrow_x = x + box_w + Inches(0.05)
            arrow_y = top + (box_h - Inches(0.5)) / 2
            add_arrow_right(slide, arrow_x, arrow_y,
                            width=arrow_w - Inches(0.1), height=Inches(0.5),
                            color=MS_DARK_BLUE)

    # Annotations below specific steps
    if annotations:
        for step_idx, text in annotations:
            ax = start_x + step_idx * (box_w + arrow_w)
            add_bullet_list(slide, text if isinstance(text, list) else [text],
                            ax, top + box_h + Inches(0.3), box_w,
                            font_size=12, color=MS_TEXT_MUTED)

    return boxes


def add_process_flow_grouped(slide, steps, group_range=None, group_label="",
                             left=CONTENT_LEFT, top=Inches(2.0),
                             box_w=Inches(2.5), box_h=Inches(1.4),
                             arrow_w=Inches(0.6), colors=None,
                             annotations=None):
    """Add a process flow with an optional dashed grouping border.

    Like add_process_flow but can draw a dashed rounded rectangle around a
    subset of steps with a label (e.g., "Wide-adoption" grouping in screenshot 4).

    group_range: (start_idx, end_idx) inclusive indices of steps to group.
    group_label: text label for the group (appears above the dashed border).

    Returns (list_of_box_shapes, group_shape_or_None).
    """
    n = len(steps)
    if colors is None:
        palette = [MS_DARK_BLUE, MS_BLUE, MS_LIGHT_BLUE]
        colors = [palette[i % len(palette)] for i in range(n)]

    total_w = n * box_w + (n - 1) * arrow_w
    start_x = left + (CONTENT_WIDTH - total_w) / 2

    group_shape = None
    # Draw grouping border FIRST (behind boxes)
    if group_range is not None:
        gs, ge = group_range
        gx = start_x + gs * (box_w + arrow_w) - Inches(0.15)
        gw = (ge - gs + 1) * box_w + (ge - gs) * arrow_w + Inches(0.3)
        g_top = top - Inches(0.25)
        g_h = box_h + Inches(0.5)
        group_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, gx, g_top, gw, g_h)
        group_shape.fill.background()
        group_shape.line.color.rgb = MS_PURPLE
        group_shape.line.width = Pt(2.0)
        group_shape.line.dash_style = 4  # dash
        try:
            group_shape.adjustments[0] = 0.06
        except Exception:
            pass
        if group_label:
            add_textbox(slide, group_label, gx, g_top + g_h + Inches(0.05),
                        gw, Inches(0.35),
                        font_size=16, color=MS_PURPLE, bold=True,
                        alignment=PP_ALIGN.CENTER, italic=True)

    # Draw boxes and arrows (on top of grouping)
    boxes = add_process_flow(slide, steps, left, top, box_w, box_h,
                             arrow_w, colors, annotations)
    return boxes, group_shape


# ═════════════════════════════════════════════════════════════
# UTILITIES
# ═════════════════════════════════════════════════════════════

def parse_slide_notes(notes_path):
    """Parse a slide markdown file and extract speaker notes by slide index (0-based).

    Returns a list of strings, one per slide. Slides without notes get an empty string.
    The notes are extracted verbatim from <!-- ... --> HTML comments in each slide.
    Usage:
        notes = parse_slide_notes("samples/slides/my-deck.md")
        create_lead_slide(prs, ..., notes=notes[0])
        create_standard_slide(prs, ..., notes=notes[2])
    """
    import re
    with open(notes_path, "r") as f:
        content = f.read()

    # Remove the YAML frontmatter + <style> block (everything before first slide content)
    # Split on slide separator
    slides = content.split("\n---\n")
    result = []
    for slide in slides:
        # Find HTML comments that contain speaker notes
        # Match both formats: <!-- note --> and <!--\nnote\n-->
        comments = re.findall(r"<!--\s*(.*?)\s*-->", slide, re.DOTALL)
        if comments:
            # Filter out directive-only comments (start with _class, _paginate, etc.)
            note_texts = [c.strip() for c in comments
                          if c.strip() and not c.strip().startswith("_")]
            result.append(note_texts[-1] if note_texts else "")
        else:
            result.append("")
    return result


def download_image(url, filepath):
    """Download an image from URL to local path."""
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=10) as resp:
            with open(filepath, "wb") as f:
                f.write(resp.read())
        return True
    except Exception as e:
        print(f"  [WARN] Could not download {url}: {e}")
        return False


def save_presentation(prs, output_path):
    """Save the presentation and print summary."""
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"\n  Presentation saved to: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")
    return output_path
